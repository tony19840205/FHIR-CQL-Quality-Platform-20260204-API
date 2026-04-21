# -*- coding: utf-8 -*-
"""
國際競爭者分析 PITCH PPT 產生器
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 色彩定義（白底列印版）=====
DARK_BG = RGBColor(0xff, 0xff, 0xff)  # 白色背景
ACCENT_BLUE = RGBColor(0x00, 0x56, 0xb3)
ACCENT_GREEN = RGBColor(0x00, 0x8a, 0x55)
ACCENT_ORANGE = RGBColor(0xe0, 0x7c, 0x00)
ACCENT_RED = RGBColor(0xcc, 0x33, 0x33)
ACCENT_PURPLE = RGBColor(0x7b, 0x43, 0x9a)
WHITE = RGBColor(0x22, 0x22, 0x22)  # 深色文字
LIGHT_GRAY = RGBColor(0x55, 0x55, 0x55)  # 中灰文字
DARK_GRAY = RGBColor(0x99, 0x99, 0x99)
CARD_BG = RGBColor(0xf0, 0xf0, 0xf5)  # 淺灰卡片
HIGHLIGHT_YELLOW = RGBColor(0xcc, 0x8a, 0x00)  # 深金色

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = "Microsoft JhengHei"
    return txBox

def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = "Microsoft JhengHei"
    p.space_before = space_before
    return p

# ===================================================================
# SLIDE 1: 封面
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "FHIR CQL 自動化醫療品質與健康監測平台", 40, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "國際競爭者分析 International Competitor Analysis", 28, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             "台北長庚紀念醫院 ｜ 2026 電子病歷應用程式展覽", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# 底部裝飾線
line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.5), Inches(5), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_BLUE
line_shape.line.fill.background()

# ===================================================================
# SLIDE 2: 全球 FHIR/CQL 競爭格局
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "全球 FHIR/CQL 競爭格局", 32, WHITE, True, PP_ALIGN.LEFT)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(4), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_BLUE
line_shape.line.fill.background()

competitors = [
    ("Smile Digital Health", "加拿大", "$44M Series B",
     "🎯 主攻：FHIR 臨床數據管理平台",
     "🏢 市值：估計 $150M-200M USD（未上市）",
     "💵 營收：估計 $15M-25M/年",
     "190+ 部署\n1.1億美國人覆蓋\n維護 HAPI FHIR\n（全球最大開源FHIR Server）", ACCENT_BLUE),
    ("b.well Connected Health", "美國", "$52.5M Series B",
     "🎯 主攻：消費者端健康資料聚合",
     "🏢 市值：估計 $200M-250M USD（未上市）",
     "💵 營收：估計 $10M-20M/年",
     "消費者健康平台\nFHIR資料聚合\n80+ 醫療系統整合\n患者端健康管理", ACCENT_GREEN),
    ("Firely", "荷蘭", "未公開",
     "🎯 主攻：FHIR 開發工具與標準化",
     "🏢 市值：估計 $30M-50M USD（未上市）",
     "💵 營收：估計 $5M-10M/年",
     "FHIR .NET SDK 維護者\nSimplifier.net 平台\n歐洲 FHIR 標準制定者\n政府級合規工具", ACCENT_ORANGE),
    ("Lantana Consulting", "美國", "未公開(政府合約)",
     "🎯 主攻：CQL/eCQM 品質標準制定",
     "🏢 市值：估計 $20M-40M USD（未上市）",
     "💵 營收：估計 $8M-15M/年（政府標案）",
     "CQL 標準核心制定者\neCQM 品質指標專家\nCMS/ONC 政府合約\n美國醫療品質基礎設施", ACCENT_PURPLE),
]

for i, (name, country, funding, focus, market_cap, revenue, desc, color) in enumerate(competitors):
    left = Inches(0.4 + i * 3.15)
    top = Inches(1.5)
    card = add_shape(slide, left, top, Inches(2.95), Inches(5.5), CARD_BG, color)
    
    # 公司名
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), Inches(2.65), Inches(0.5),
                 name, 18, color, True)
    # 國家
    add_text_box(slide, left + Inches(0.15), top + Inches(0.6), Inches(2.65), Inches(0.3),
                 f"📍 {country}", 13, LIGHT_GRAY)
    # 融資
    add_text_box(slide, left + Inches(0.15), top + Inches(0.9), Inches(2.65), Inches(0.3),
                 f"💰 {funding}", 13, HIGHLIGHT_YELLOW, True)
    # 主攻/市值/營收
    txBox_info = add_text_box(slide, left + Inches(0.15), top + Inches(1.25), Inches(2.65), Inches(1.2),
                              "", 11, WHITE)
    tf_info = txBox_info.text_frame
    tf_info.paragraphs[0].text = focus
    tf_info.paragraphs[0].font.size = Pt(11)
    tf_info.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x56, 0xb3)
    tf_info.paragraphs[0].font.name = "Microsoft JhengHei"
    add_paragraph(tf_info, market_cap, 11, RGBColor(0x00, 0x56, 0xb3), space_before=Pt(2))
    add_paragraph(tf_info, revenue, 11, RGBColor(0x00, 0x56, 0xb3), space_before=Pt(2))
    # 描述
    txBox = add_text_box(slide, left + Inches(0.15), top + Inches(2.6), Inches(2.65), Inches(2.7),
                         "", 13, WHITE)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split("\n"):
        add_paragraph(tf, f"• {line}", 13, WHITE, space_before=Pt(4))

# ===================================================================
# SLIDE 3: Smile Digital Health 深度分析
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "最大競爭者：Smile Digital Health", 32, ACCENT_BLUE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(5), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_BLUE
line_shape.line.fill.background()

# 左：公司概況
card1 = add_shape(slide, Inches(0.4), Inches(1.4), Inches(6), Inches(5.5), CARD_BG)
txBox = add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.6), Inches(0.4),
                     "公司概況", 20, ACCENT_BLUE, True)
tf = txBox.text_frame

items = [
    "成立：2016年，加拿大多倫多",
    "產品：OmniQ™ 臨床數據管理平台",
    "融資：Series B $44M USD（2022）",
    "規模：200+ 員工",
    "客戶：190+ 醫療機構部署",
    "覆蓋：1.1 億美國人",
    "開源貢獻：維護 HAPI FHIR Server",
    "  → GitHub 2,300+ Stars, 274 Contributors",
    "認證：NCQA 驗證、HL7 FHIR 合規"
]
for item in items:
    add_paragraph(tf, f"• {item}", 14, WHITE, space_before=Pt(3))

# 右：產品特色
card2 = add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5), CARD_BG)
txBox2 = add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4),
                      "OmniQ™ 核心功能", 20, ACCENT_GREEN, True)
tf2 = txBox2.text_frame

features = [
    "FHIR R4 原生架構",
    "臨床數據倉儲 (CDR)",
    "CDS Hooks 臨床決策支援",
    "SMART on FHIR 應用生態系",
    "患者同意管理",
    "跨機構資料交換",
    "API 管理與安全閘道",
    "合規報告（HEDIS/eCQM）",
]
for f in features:
    add_paragraph(tf2, f"✦ {f}", 14, WHITE, space_before=Pt(4))

add_paragraph(tf2, "", 10, WHITE, space_before=Pt(10))
add_paragraph(tf2, "年營收估計：$15M-25M USD", 16, HIGHLIGHT_YELLOW, True, space_before=Pt(6))

# ===================================================================
# SLIDE 4: 競爭優劣勢比較表
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "功能比較矩陣", 32, WHITE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(3), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_BLUE
line_shape.line.fill.background()

# 表格
headers = ["功能", "台北長庚\nFHIR CQL 平台", "Smile\nDigital Health", "b.well\nConnected", "Firely", "Lantana\nConsulting"]
col_widths = [Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]

rows_data = [
    ["FHIR R4 標準", "✅", "✅", "✅", "✅", "✅"],
    ["CQL 品質指標引擎", "✅ 50指標", "⚠️ 有限", "❌", "❌", "✅ 制定者"],
    ["即時監測儀表板", "✅ 4大畫面", "✅", "✅", "❌", "❌"],
    ["民眾健康入口", "✅", "⚠️", "✅ 核心", "❌", "❌"],
    ["跨院資料交換", "✅", "✅", "✅", "✅", "⚠️"],
    ["台灣健保適配", "✅ 原生", "❌", "❌", "❌", "❌"],
    ["開源/透明度", "✅ GitHub", "✅ HAPI", "❌", "✅ SDK", "⚠️"],
    ["部署成本", "🟢 極低", "🔴 企業級", "🔴 企業級", "🟡 中等", "🔴 政府標"],
    ["在地化程度", "🟢 台灣原生", "🔴 北美", "🔴 北美", "🟡 歐洲", "🔴 北美"],
]

table_shape = slide.shapes.add_table(len(rows_data) + 1, 6,
                                      Inches(0.4), Inches(1.4),
                                      Inches(12.5), Inches(5.5))
table = table_shape.table

# Style header
for j, header in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        paragraph.font.name = "Microsoft JhengHei"
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE if j == 0 else RGBColor(0x34, 0x4e, 0x6a)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# 第二欄（台北長庚）header 特殊顏色
cell = table.cell(0, 1)
cell.fill.solid()
cell.fill.fore_color.rgb = ACCENT_GREEN

# Fill data
for i, row in enumerate(rows_data):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            paragraph.font.name = "Microsoft JhengHei"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        if j == 1:  # 台北長庚列高亮
            cell.fill.fore_color.rgb = RGBColor(0xd5, 0xf0, 0xe0)
        else:
            cell.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xfa) if i % 2 == 0 else RGBColor(0xe8, 0xe8, 0xf0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# ===================================================================
# SLIDE 5: 我們的獨特優勢
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "我們的不可取代優勢", 32, WHITE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(4), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_GREEN
line_shape.line.fill.background()

advantages = [
    ("🏥", "台灣健保原生適配", 
     "國際競爭者皆無台灣健保制度經驗\n我們的 50 項指標完全對應台灣醫療品質規範\n無需二次開發，即插即用",
     ACCENT_BLUE),
    ("💰", "極低部署成本",
     "Smile 等企業級方案年費 $50K-500K USD\n我們基於開源技術棧、Render 免費部署\n單人即可維護，適合中小型醫療機構",
     ACCENT_GREEN),
    ("⚡", "CQL 50 指標全自動化",
     "全場唯一完整實作 CQL 引擎的品質平台\n涵蓋 6 大醫療領域、50 項指標\n即時計算、無須人工審查",
     ACCENT_ORANGE),
    ("🏆", "競賽實證",
     "50 個入選作品中排名 #1\n同一人同時入選 2 個作品 (#1 & #5)\n技術深度獲評審認可",
     ACCENT_RED),
]

for i, (icon, title, desc, color) in enumerate(advantages):
    left = Inches(0.4 + (i % 2) * 6.3)
    top = Inches(1.5 + (i // 2) * 2.8)
    card = add_shape(slide, left, top, Inches(6.0), Inches(2.5), CARD_BG, color)
    
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(5.5), Inches(0.5),
                 f"{icon} {title}", 22, color, True)
    
    txBox = add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(5.5), Inches(1.6),
                         "", 14, WHITE)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split("\n"):
        add_paragraph(tf, line, 14, LIGHT_GRAY, space_before=Pt(3))

# ===================================================================
# SLIDE 6: 市場定位地圖
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "市場定位：成本 vs 功能完整度", 32, WHITE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(4), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_BLUE
line_shape.line.fill.background()

# 座標軸
# Y軸 (成本)
y_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5), Pt(2), Inches(5))
y_axis.fill.solid()
y_axis.fill.fore_color.rgb = LIGHT_GRAY
y_axis.line.fill.background()

add_text_box(slide, Inches(0.3), Inches(1.5), Inches(1.5), Inches(0.4),
             "高成本 ↑", 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)
add_text_box(slide, Inches(0.3), Inches(6.0), Inches(1.5), Inches(0.4),
             "低成本 ↓", 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# X軸 (功能)
x_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.0), Inches(10), Pt(2))
x_axis.fill.solid()
x_axis.fill.fore_color.rgb = LIGHT_GRAY
x_axis.line.fill.background()

add_text_box(slide, Inches(2), Inches(6.5), Inches(2), Inches(0.4),
             "← 功能有限", 12, LIGHT_GRAY)
add_text_box(slide, Inches(10), Inches(6.5), Inches(2), Inches(0.4),
             "功能完整 →", 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# 定位點
positions = [
    ("Smile Digital Health", "估值$150-200M", Inches(9), Inches(2.0), ACCENT_BLUE),
    ("b.well", "估值$200-250M", Inches(7), Inches(2.5), ACCENT_GREEN),
    ("Firely", "估值$30-50M", Inches(5), Inches(3.0), ACCENT_ORANGE),
    ("Lantana", "估值$20-40M", Inches(8), Inches(2.8), ACCENT_PURPLE),
    ("台北長庚\nFHIR CQL 平台", "$0 融資", Inches(9.5), Inches(5.5), ACCENT_RED),
]

for name, valuation, x, y, color in positions:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    
    is_us = "台北" in name
    add_text_box(slide, x - Inches(0.5), y + Inches(0.5), Inches(1.8), Inches(0.4),
                 name, 11 if not is_us else 13, 
                 HIGHLIGHT_YELLOW if is_us else color,
                 bold=is_us,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x - Inches(0.5), y + Inches(0.85), Inches(1.8), Inches(0.3),
                 valuation, 9,
                 HIGHLIGHT_YELLOW if is_us else LIGHT_GRAY,
                 bold=is_us,
                 alignment=PP_ALIGN.CENTER)

# 標註
highlight_box = add_shape(slide, Inches(7.5), Inches(4.8), Inches(4.5), Inches(1.0), 
                           RGBColor(0xd5, 0xf0, 0xe0), ACCENT_GREEN)
add_text_box(slide, Inches(7.7), Inches(4.9), Inches(4.1), Inches(0.8),
             "⭐ 最佳甜蜜點：功能完整度媲美國際大廠\n    但部署成本僅為其 1/100", 14, RGBColor(0x00, 0x5a, 0x35), True)

# ===================================================================
# SLIDE 7: 與 Lantana 的 CQL 深度比較
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "CQL 引擎深度比較：vs Lantana Consulting Group", 28, WHITE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(5), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_PURPLE
line_shape.line.fill.background()

# Lantana
card_l = add_shape(slide, Inches(0.4), Inches(1.4), Inches(6), Inches(5.5), CARD_BG, ACCENT_PURPLE)
txBox = add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.6), Inches(0.4),
                     "Lantana Consulting Group", 20, ACCENT_PURPLE, True)
tf = txBox.text_frame
items_l = [
    "CQL 標準的核心制定者之一",
    "eCQM (electronic Clinical Quality Measures) 專家",
    "美國 CMS / ONC 政府合約承包商",
    "負責美國醫療品質報告基礎設施",
    "客群：美國聯邦政府、大型保險公司",
    "定價：政府標案等級（$M/年）",
    "",
    "⚠ 限制：",
    "  • 僅服務美國市場",
    "  • 指標僅覆蓋 HEDIS/eCQM 美規",
    "  • 無台灣健保適配",
    "  • 非即時監測，報告導向",
]
for item in items_l:
    add_paragraph(tf, item, 13, WHITE if not item.startswith("⚠") else ACCENT_ORANGE,
                  bold=item.startswith("⚠"), space_before=Pt(3))

# 台北長庚
card_r = add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5), CARD_BG, ACCENT_GREEN)
txBox2 = add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4),
                      "台北長庚 FHIR CQL 平台", 20, ACCENT_GREEN, True)
tf2 = txBox2.text_frame
items_r = [
    "CQL 引擎在地化完整實作",
    "50 項台灣醫療品質指標",
    "涵蓋 6 大領域：",
    "  🫀 心血管 / 🩺 糖尿病 / 🧠 精神科",
    "  💊 抗生素 / 🏥 急診 / 🔬 感控",
    "即時計算 + 視覺化儀表板",
    "民眾端健康入口網站",
    "部署成本：接近零",
    "",
    "✅ 優勢：",
    "  • 台灣唯一完整 CQL 品質平台",
    "  • 即時監測 vs 報告導向",
    "  • 開源透明，GitHub 可驗證",
]
for item in items_r:
    add_paragraph(tf2, item, 13, WHITE if not item.startswith("✅") else ACCENT_GREEN,
                  bold=item.startswith("✅"), space_before=Pt(3))

# ===================================================================
# SLIDE 8: 國際專利分析
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "FHIR/CQL 國際專利佈局分析", 32, WHITE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(4), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_BLUE
line_shape.line.fill.background()

patents = [
    ("Smile Digital Health", "加拿大/美國", [
        "Google Patents 查無任何專利紀錄",
        "",
        "策略：走純開源路線",
        "維護 HAPI FHIR (Apache 2.0)",
        "以社群影響力為護城河",
        "不以專利作為商業壁壘",
    ], "0 件 ❌", ACCENT_BLUE),
    ("b.well Connected Health", "美國", [
        "US11,004,547B2（已核准 2021）",
        "— 多資料來源健康數據聚合系統",
        "US20210313071A1（申請中）",
        "— 健康風險動態評估",
        "",
        "優先權日：2016-11-01",
        "發明人：Kristen Valdes (CEO)",
        "方向：消費者端個人健康管理",
    ], "4 件（同一專利族）", ACCENT_GREEN),
    ("Firely / HL7", "荷蘭/國際", [
        "Google Patents 查無 FHIR 相關專利",
        "",
        "策略：純開源 SDK",
        "FHIR .NET SDK (BSD 授權)",
        "透過 HL7 標準組織影響力",
        "不依賴專利保護",
    ], "0 件 ❌", ACCENT_ORANGE),
    ("Lantana Consulting", "美國", [
        "Google Patents 查無任何專利紀錄",
        "",
        "策略：政府合約導向",
        "CQL 為 HL7 開放標準",
        "核心價值在標準制定話語權",
        "收入來自 CMS/ONC 政府標案",
    ], "0 件 ❌", ACCENT_PURPLE),
]

for i, (company, country, patent_list, count, color) in enumerate(patents):
    left = Inches(0.4 + i * 3.15)
    top = Inches(1.4)
    card = add_shape(slide, left, top, Inches(2.95), Inches(5.5), CARD_BG, color)
    
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(2.65), Inches(0.4),
                 company, 16, color, True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5), Inches(2.65), Inches(0.3),
                 f"📍 {country}", 12, LIGHT_GRAY)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.8), Inches(2.65), Inches(0.3),
                 f"📋 專利數量：{count}", 12, HIGHLIGHT_YELLOW, True)
    
    txBox = add_text_box(slide, left + Inches(0.15), top + Inches(1.3), Inches(2.65), Inches(3.8),
                         "", 11, WHITE)
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    for pat in patent_list:
        add_paragraph(tf, f"• {pat}", 11, LIGHT_GRAY, space_before=Pt(4))

# ===================================================================
# SLIDE 9: 我們的專利
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "我們的專利佈局", 32, WHITE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(3), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_GREEN
line_shape.line.fill.background()

# 大卡片 - 專利資訊
patent_card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.3), CARD_BG, ACCENT_GREEN)

# 專利標題
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(7), Inches(0.5),
             "🏛️ 中華民國發明專利", 24, ACCENT_GREEN, True)

# 專利號
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(7), Inches(0.4),
             "專利案號：114149490", 18, HIGHLIGHT_YELLOW, True)

# 詳細資訊
info_items = [
    ("中文名稱", "多伺服器整合與資料分析系統與方法"),
    ("英文名稱", "System and Method for Multi-server\n                    Integration and Data Analysis"),
    ("申請人1", "長庚醫療財團法人台北長庚紀念醫院"),
    ("申請人2", "長庚醫療財團法人林口長庚紀念醫院"),
    ("發明人", "陳孟琪 CHEN, MENG-CHI"),
    ("申請日期", "114年12月16日 (2025/12/16)"),
    ("案件編號", "25044CGH-TW"),
]

y_pos = 2.8
for label, value in info_items:
    add_text_box(slide, Inches(1.0), Inches(y_pos), Inches(2.0), Inches(0.35),
                 f"【{label}】", 14, ACCENT_BLUE, True)
    add_text_box(slide, Inches(3.0), Inches(y_pos), Inches(4.8), Inches(0.45),
                 value, 14, WHITE)
    y_pos += 0.42

# 右側 - 專利價值分析
value_card = add_shape(slide, Inches(8.3), Inches(1.5), Inches(4.7), Inches(5.3), CARD_BG, ACCENT_BLUE)

add_text_box(slide, Inches(8.5), Inches(1.6), Inches(4.3), Inches(0.5),
             "⚡ 專利核心技術價值", 20, ACCENT_BLUE, True)

value_items = [
    ("多伺服器整合", "跨院區 FHIR Server 資料\n整合與同步技術", ACCENT_GREEN),
    ("資料分析引擎", "CQL 品質指標自動化\n計算與即時監測", ACCENT_ORANGE),
    ("系統架構創新", "低成本高效能的\n分散式醫療數據平台", ACCENT_PURPLE),
]

y_val = 2.3
for title, desc, color in value_items:
    val_box = add_shape(slide, Inches(8.5), Inches(y_val), Inches(4.3), Inches(1.35), 
                         RGBColor(0xf0, 0xf0, 0xf8), color)
    add_text_box(slide, Inches(8.7), Inches(y_val + 0.1), Inches(3.9), Inches(0.35),
                 f"🔹 {title}", 15, color, True)
    add_text_box(slide, Inches(8.7), Inches(y_val + 0.5), Inches(3.9), Inches(0.7),
                 desc, 13, LIGHT_GRAY)
    y_val += 1.5

# 底部亮點
highlight = add_shape(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5),
                       RGBColor(0xd5, 0xf0, 0xe0), ACCENT_GREEN)
add_text_box(slide, Inches(0.7), Inches(6.92), Inches(11.9), Inches(0.45),
             "✅ 台灣首件結合 FHIR 多伺服器整合 + CQL 品質分析的發明專利  |  雙院區聯名申請（台北長庚 + 林口長庚）",
             14, RGBColor(0x00, 0x5a, 0x35), True, PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 10: 專利定位圖（類似市場定位圖）
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "專利佈局定位：IP 保護 vs 技術完整度", 32, WHITE, True)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(5), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_GREEN
line_shape.line.fill.background()

# Y軸 (IP保護強度)
y_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5), Pt(2), Inches(5))
y_axis.fill.solid()
y_axis.fill.fore_color.rgb = LIGHT_GRAY
y_axis.line.fill.background()

add_text_box(slide, Inches(0.1), Inches(1.5), Inches(1.7), Inches(0.4),
             "強 IP 保護 ↑", 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)
add_text_box(slide, Inches(0.1), Inches(6.0), Inches(1.7), Inches(0.4),
             "無 IP 保護 ↓", 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# X軸 (技術完整度)
x_axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.0), Inches(10), Pt(2))
x_axis.fill.solid()
x_axis.fill.fore_color.rgb = LIGHT_GRAY
x_axis.line.fill.background()

add_text_box(slide, Inches(2), Inches(6.5), Inches(2.5), Inches(0.4),
             "← FHIR/CQL 功能有限", 12, LIGHT_GRAY)
add_text_box(slide, Inches(9.5), Inches(6.5), Inches(2.5), Inches(0.4),
             "FHIR+CQL 完整 →", 12, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# --- 定位點 ---
# Smile: 功能完整但無IP => 右下
dot_smile = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(5.0), Inches(0.5), Inches(0.5))
dot_smile.fill.solid()
dot_smile.fill.fore_color.rgb = ACCENT_BLUE
dot_smile.line.fill.background()
add_text_box(slide, Inches(9.0), Inches(5.5), Inches(2.0), Inches(0.35),
             "Smile Digital Health", 11, ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.0), Inches(5.8), Inches(2.0), Inches(0.3),
             "0 件專利 / 純開源", 9, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# b.well: 消費者端，不含CQL => 中間偏左，有IP
dot_bwell = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(2.8), Inches(0.5), Inches(0.5))
dot_bwell.fill.solid()
dot_bwell.fill.fore_color.rgb = ACCENT_GREEN
dot_bwell.line.fill.background()
add_text_box(slide, Inches(5.0), Inches(3.3), Inches(1.8), Inches(0.35),
             "b.well", 11, ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(3.6), Inches(2.8), Inches(0.3),
             "4 件專利 / 消費者端聚合", 9, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Firely: SDK工具，無IP => 左下
dot_firely = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.0), Inches(5.2), Inches(0.5), Inches(0.5))
dot_firely.fill.solid()
dot_firely.fill.fore_color.rgb = ACCENT_ORANGE
dot_firely.line.fill.background()
add_text_box(slide, Inches(3.5), Inches(5.7), Inches(1.8), Inches(0.35),
             "Firely", 11, ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.2), Inches(6.0), Inches(2.3), Inches(0.3),
             "0 件專利 / 純開源 SDK", 9, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Lantana: CQL標準制定但無IP => 右側中間偏下
dot_lantana = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), Inches(5.3), Inches(0.5), Inches(0.5))
dot_lantana.fill.solid()
dot_lantana.fill.fore_color.rgb = ACCENT_PURPLE
dot_lantana.line.fill.background()
add_text_box(slide, Inches(7.5), Inches(5.8), Inches(1.8), Inches(0.35),
             "Lantana", 11, ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.0), Inches(6.1), Inches(2.8), Inches(0.3),
             "0 件專利 / 政府合約", 9, LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 台北長庚: FHIR+CQL完整 + 有專利 => 右上！唯一
dot_us = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(1.8), Inches(0.6), Inches(0.6))
dot_us.fill.solid()
dot_us.fill.fore_color.rgb = ACCENT_RED
dot_us.line.fill.background()
# 星星
star = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(1.7), Inches(0.3), Inches(0.3))
star.fill.solid()
star.fill.fore_color.rgb = HIGHLIGHT_YELLOW
star.line.fill.background()
add_text_box(slide, Inches(9.2), Inches(2.5), Inches(2.5), Inches(0.4),
             "台北長庚 FHIR CQL", 14, HIGHLIGHT_YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.2), Inches(2.85), Inches(2.5), Inches(0.3),
             "發明專利 10 項請求項", 11, HIGHLIGHT_YELLOW, alignment=PP_ALIGN.CENTER)

# 右上角高亮框
highlight_box = add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(2.0),
                           RGBColor(0xd5, 0xf0, 0xe0), ACCENT_GREEN)
# 重新加文字在框上
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5.4), Inches(1.8),
             "⭐ 全球唯一甜蜜點\n\n同時具備：\n✅ FHIR + CQL 完整技術實作\n✅ 發明專利保護（10 項請求項）\n✅ 其他 4 家競爭者皆無此組合",
             13, RGBColor(0x00, 0x5a, 0x35), True)

# 資料來源標註
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
             "資料來源：Google Patents 實際查證 (2026/04/11)  |  Smile、Firely、Lantana 查無專利紀錄  |  b.well 4 件皆為消費者端健康聚合",
             9, RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 11: 結語
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             "我們站在全球 FHIR/CQL 的前沿", 36, WHITE, True, PP_ALIGN.CENTER)

# 三個關鍵數字
stats = [
    ("$96.5M+", "國際競爭者\n累計融資總額", ACCENT_BLUE),
    ("$0", "我們的\n外部融資", ACCENT_GREEN),
    ("50 項", "CQL 品質指標\n台灣唯一完整實作", ACCENT_ORANGE),
]

for i, (num, desc, color) in enumerate(stats):
    left = Inches(1.5 + i * 3.8)
    card = add_shape(slide, left, Inches(3.0), Inches(3.2), Inches(2.5), CARD_BG, color)
    add_text_box(slide, left, Inches(3.2), Inches(3.2), Inches(1.0),
                 num, 36, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(4.3), Inches(3.2), Inches(0.8),
                 desc, 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
             "「不是我們缺乏資源，是我們用更少的資源，做到了同等的事。」", 
             22, HIGHLIGHT_YELLOW, True, PP_ALIGN.CENTER)

# ===== 儲存 =====
output_path = os.path.join(os.path.dirname(__file__), "FHIR_CQL_國際競爭者分析_PITCH_v4.pptx")
prs.save(output_path)
print(f"PPT 已儲存至: {output_path}")
