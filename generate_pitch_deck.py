# -*- coding: utf-8 -*-
"""
FHIR-CQL Quality Platform - Investor Pitch Deck Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x0F, 0x2B, 0x46)
BRAND_BLUE = RGBColor(0x00, 0x6D, 0xAA)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x00, 0xA6, 0x5A)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
RED = RGBColor(0xE8, 0x3E, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper Functions ──

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft JhengHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_name='Microsoft JhengHei'):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(6)
    return txBox

def add_table_slide(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=BRAND_BLUE):
    """data: list of lists, first row is header"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = 'Microsoft JhengHei'
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = BLACK
                    paragraph.font.bold = False
            
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

def add_card(slide, left, top, width, height, title, content, title_color=BRAND_BLUE, bg_color=WHITE):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)
    
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    
    # Content
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                 content, font_size=12, color=GRAY, bold=False)

def slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 14

# ══════════════════════════════════════════════
# SLIDE 1: Cover
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

# Accent bar
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_ORANGE)

# Title block
add_multi_text(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(5), [
    ("FHIR-CQL 醫療品質量測平台", 44, WHITE, True, PP_ALIGN.LEFT),
    ("", 12, WHITE, False, PP_ALIGN.LEFT),
    ("台灣首個 FHIR 原生品質指標自動計算引擎", 28, ACCENT_BLUE, False, PP_ALIGN.LEFT),
    ("", 12, WHITE, False, PP_ALIGN.LEFT),
    ("39 項品質指標 ｜ 傳染病監控 ｜ 公衛追蹤 ｜ ESG 永續報告", 20, GRAY, False, PP_ALIGN.LEFT),
    ("", 24, WHITE, False, PP_ALIGN.LEFT),
    ("種子輪募資簡報  |  2026 年 4 月", 18, GRAY, False, PP_ALIGN.LEFT),
])

# Badge
badge = add_shape_bg(slide, Inches(9.5), Inches(5.5), Inches(3.2), Inches(1.2), ACCENT_ORANGE)
add_text_box(slide, Inches(9.6), Inches(5.55), Inches(3.0), Inches(0.5),
             "🏆 衛福部 SMART 50 得獎", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.6), Inches(6.05), Inches(3.0), Inches(0.5),
             "📋 專利申請中", font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

slide_number(slide, 1, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 2: Problem — 痛點
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "❌  問題：醫院品質指標的困境", font_size=32, color=WHITE, bold=True)

# Pain point cards
cards = [
    ("⏱️ 人工統計", "全台 400+ 家醫院，每季手動\n統計數十項品質指標\n\n每院投入 2-3 名專責人力\n全國每年浪費數億人力成本"),
    ("🔀 系統不相容", "各醫院 HIS 系統皆不相同\n傳統方式：每院客製 SQL\n接 400 家要做 400 次\n\n擴展成本與醫院數成正比"),
    ("📊 數據孤島", "各醫院指標數據各自獨立\n無法跨院比較\n無法即時監控全國趨勢\n\n寶貴數據被鎖在各院圍牆內"),
    ("📋 政策壓力", "衛福部次世代電子病歷計畫\n強制推行 FHIR 標準\n品質指標即將要求電子化申報\n\n醫院急需解決方案"),
]

for i, (title, content) in enumerate(cards):
    x = Inches(0.5 + i * 3.15)
    # Card bg
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.9), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 0 else LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    card.line.width = Pt(1)
    
    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.6),
                 title, font_size=20, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(3.5),
                 content, font_size=15, color=BLACK, bold=False)

# Bottom highlight
add_shape_bg(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
             "Health Catalyst 用傳統架構做了 18 年，每院導入要 6-12 個月、50-100 萬美元，市值已從 60 億跌到 8 億", 
             font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(slide, 2, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 3: Solution — 解決方案
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BRAND_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "✅  解決方案：FHIR + CQL 自動品質量測引擎", font_size=32, color=WHITE, bold=True)

# USB-C analogy
add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5), LIGHT_BLUE)
add_multi_text(slide, Inches(0.8), Inches(1.7), Inches(5.3), Inches(2.3), [
    ("🔌 USB-C 比喻", 22, DARK_BLUE, True, PP_ALIGN.LEFT),
    ("", 8, BLACK, False, PP_ALIGN.LEFT),
    ("傳統方式 = 每家醫院不同插座，做 400 個轉接頭", 16, RED, False, PP_ALIGN.LEFT),
    ("我們的方式 = USB-C 統一介面，插上就能用", 16, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("", 8, BLACK, False, PP_ALIGN.LEFT),
    ("衛福部已規定所有醫院要裝 FHIR（USB-C）", 15, GRAY, False, PP_ALIGN.LEFT),
    ("我們提前做好了充電器", 15, BRAND_BLUE, True, PP_ALIGN.LEFT),
])

# Key metrics
metrics = [
    ("39", "項品質指標"),
    ("50+", "CQL 查詢檔"),
    ("645+", "測試病人"),
    ("4", "大儀表板"),
]
for i, (num, label) in enumerate(metrics):
    x = Inches(7.0 + i * 1.5)
    add_text_box(slide, x, Inches(1.8), Inches(1.4), Inches(0.8),
                 num, font_size=36, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.5), Inches(1.4), Inches(0.5),
                 label, font_size=14, color=GRAY, bold=False, alignment=PP_ALIGN.CENTER)

# Comparison table
data = [
    ["", "傳統資料庫方式", "FHIR + CQL（我們）"],
    ["接入新醫院", "3-6 個月/院", "1-2 週/院"],
    ["工程師需求", "每院需了解 DB schema", "CQL 通用，不需客製"],
    ["指標邏輯修改", "改 400 套 SQL", "改 1 份 CQL，全部生效"],
    ["跨院比較", "欄位不同，耗人力對齊", "FHIR 標準化，天生可比"],
    ["邊際成本", "O(n) 隨醫院數成長", "O(1) 趨近於零"],
]
add_table_slide(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.8),
                6, 3, data,
                col_widths=[Inches(2.5), Inches(4.9), Inches(4.9)])

slide_number(slide, 3, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 4: Product — 四大儀表板
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "🖥️  產品：四大即時儀表板", font_size=32, color=WHITE, bold=True)

dashboards = [
    ("📊 醫療品質指標", "39 項健保品質指標\n自動計算分子/分母\n涵蓋用藥安全、門診、\n住院、手術、結果面", BRAND_BLUE),
    ("🦠 傳染病監控", "COVID-19、流感、登革熱\n腸病毒、急性腹瀉\nICD-10 + SNOMED CT\n即時疫情趨勢追蹤", RED),
    ("🏥 公共衛生", "COVID-19 疫苗覆蓋率\n流感疫苗接種率\n高血壓活動個案追蹤\nWHO 標準對接", ACCENT_GREEN),
    ("🌱 ESG 永續指標", "抗生素 AWaRe 分級\n電子病歷採用率\n醫療廢棄物管理\nGRI 306 標準", RGBColor(0x8B, 0x5E, 0x3C)),
]

for i, (title, content, color) in enumerate(dashboards):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.6)
    
    # Top color bar
    add_shape_bg(slide, x, y, Inches(3.0), Inches(0.08), color)
    
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.08), Inches(3.0), Inches(4.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1.5)
    
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.6), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.6), Inches(3.0),
                 content, font_size=15, color=BLACK)

# Tech stack bar
add_shape_bg(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.9), LIGHT_GRAY)
add_text_box(slide, Inches(0.6), Inches(6.25), Inches(12), Inches(0.8),
             "技術棧：FHIR R4  |  CQL (Clinical Quality Language)  |  TW Core  |  ICD-10 / SNOMED CT / LOINC  |  Express.js  |  SMART on FHIR OAuth 2.0",
             font_size=14, color=GRAY, bold=False, alignment=PP_ALIGN.CENTER)

slide_number(slide, 4, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 5: 品質指標細節
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BRAND_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "📋  39 項品質指標全覽", font_size=32, color=WHITE, bold=True)

data = [
    ["類別", "數量", "代表性指標"],
    ["用藥安全", "2", "門診注射率、門診抗生素使用率"],
    ["藥物重疊", "16", "同院/跨院重複用藥（降壓、降血脂、降血糖等 8 類）"],
    ["門診品質", "5", "慢性病連續處方、多重用藥、兒童氣喘急診、糖尿病 HbA1c"],
    ["住院品質", "6", "14 天非計畫再住院、3 天內急診、剖腹產率（4 類）"],
    ["手術品質", "8", "清潔手術抗生素、體外震波碎石、人工膝關節感染率"],
    ["結果面", "2", "AMI 死亡率、失智症安寧療護利用率"],
]
add_table_slide(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.5),
                7, 3, data,
                col_widths=[Inches(2.0), Inches(1.3), Inches(9.0)])

add_multi_text(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.8), [
    ("每項指標都有完整的 CQL 定義 → 可在任何 FHIR Server 上即時計算分子、分母、百分比", 16, DARK_BLUE, True, PP_ALIGN.LEFT),
    ("對接台灣健保署品質指標代碼，可直接用於醫院評鑑申報", 15, GRAY, False, PP_ALIGN.LEFT),
    ("支援多重編碼系統：ICD-10-CM / SNOMED CT / LOINC / NHI / ATC 藥物分類", 15, GRAY, False, PP_ALIGN.LEFT),
])

slide_number(slide, 5, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 6: Business Model — 商業模式
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "💰  商業模式：免費工具 → 數據聚合 → 數據變現", font_size=32, color=WHITE, bold=True)

# Three acts
acts = [
    ("第一幕：搶覆蓋率", "🆓", "免費幫醫院算品質指標\n↓\n醫院省下人力成本\n我們取得聚合後的指標數據\n（去識別化統計值，非個資）\n\n目標：2 年內 100+ 家醫院", 
     ACCENT_ORANGE, "現在 → Year 2"),
    ("第二幕：數據平台", "📊", "全台即時醫療品質數據集\n↓\n藥廠：疾病指標趨勢\n醫材商：手術品質數據\n保險公司：風險定價\n管顧機構：跨院比較報告\n\n「台灣版品質指標 IQVIA」", 
     BRAND_BLUE, "Year 2 → Year 4"),
    ("第三幕：國際擴展", "🌏", "CQL 是國際標準\n換 ValueSet 即可適配各國\n↓\n東南亞 FHIR 基礎建設外銷\n台灣 FHIR Server + 指標引擎\n打包輸出\n\n搭政府智慧醫療南向政策", 
     ACCENT_GREEN, "Year 4+"),
]

for i, (title, icon, content, color, timeline) in enumerate(acts):
    x = Inches(0.4 + i * 4.2)
    
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.9), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.4), Inches(1.8), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, x + Inches(1.4), Inches(1.85), Inches(0.9), Inches(0.8),
                 icon, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.2), Inches(2.85), Inches(3.5), Inches(0.4),
                 title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(3.0),
                 content, font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)
    
    # Timeline badge
    badge = add_shape_bg(slide, x + Inches(0.8), Inches(6.4), Inches(2.3), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.8), Inches(6.4), Inches(2.3), Inches(0.35),
                 timeline, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Arrow between cards
    if i < 2:
        add_text_box(slide, Inches(4.1 + i * 4.2), Inches(3.8), Inches(0.5), Inches(0.5),
                     "→", font_size=30, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(slide, 6, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 7: Data Value — 數據價值
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BRAND_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "📈  數據價值：誰會買？花多少錢？", font_size=32, color=WHITE, bold=True)

data = [
    ["客戶類型", "他們要什麼", "價值等級"],
    ["藥廠", "特定疾病指標趨勢（如：糖尿病 HbA1c 檢測率變化）→ 評估藥品市場", "💰💰💰 高"],
    ["醫材公司", "手術品質指標（如：人工膝關節感染率）→ 證明產品效果", "💰💰💰 中高"],
    ["保險公司", "醫院品質排名 → 風險定價依據", "💰💰💰 高"],
    ["管顧/研究機構", "跨院比較分析報告", "💰💰 中"],
    ["政府/NGO", "公衛政策效果評估", "💰 中低但量大"],
]
add_table_slide(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.0),
                6, 3, data,
                col_widths=[Inches(2.5), Inches(7.3), Inches(2.5)])

# Value formula
add_shape_bg(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.0), LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.05), Inches(11.7), Inches(0.9),
             "數據價值公式：   醫院數量  ×  指標種類  ×  時間序列長度  =  壟斷性資產",
             font_size=22, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Reference companies
add_multi_text(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.2), [
    ("參考案例：IQVIA 市值 400 億美元（聚合處方數據）| Flatiron Health 被 Roche 以 19 億美元收購（聚合腫瘤數據）", 15, GRAY, False, PP_ALIGN.CENTER),
    ("我們的定位：台灣版「品質指標 IQVIA」— 規模小但模型相同，護城河一旦建立不可撼動", 15, BRAND_BLUE, True, PP_ALIGN.CENTER),
])

slide_number(slide, 7, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 8: Market — 市場規模
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "🌐  市場規模", font_size=32, color=WHITE, bold=True)

# TAM / SAM / SOM
circles_data = [
    ("TAM 全球", "全球醫療品質分析市場\n~150 億美元\n（含 Health Catalyst、Premier、\nVizient 等同類產品市場）", 
     Inches(4.0), Inches(4.0), RGBColor(0xCC, 0xDD, 0xFF), Inches(0.5), Inches(1.8)),
    ("SAM 亞太", "亞太 FHIR 品質量測市場\n~5 億美元\n（台灣 + 日韓 + 東南亞\n正在推 FHIR 的國家）", 
     Inches(3.0), Inches(3.0), RGBColor(0x99, 0xBB, 0xFF), Inches(5.5), Inches(2.3)),
    ("SOM 台灣", "台灣醫院品質管理\n~3-5 億台幣\n（400+ 醫院 ×\n品質管理年預算）", 
     Inches(2.2), Inches(2.2), BRAND_BLUE, Inches(9.5), Inches(2.8)),
]

for label, text, w, h, color, x, y in circles_data:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.color.rgb = BRAND_BLUE
    circle.line.width = Pt(2)
    
    is_som = label == "SOM 台灣"
    txt_color = WHITE if is_som else DARK_BLUE
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.4),
                 label, font_size=16, color=txt_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), w - Inches(0.3), h - Inches(0.8),
                 text, font_size=11, color=txt_color, alignment=PP_ALIGN.CENTER)

# Bottom timeline
add_shape_bg(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0), LIGHT_GRAY)
add_multi_text(slide, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.9), [
    ("短期：全台 400+ 醫院品質管理部門（SaaS 訂閱）  →  中期：衛福部 TW Health App Space 上架  →  長期：東南亞 FHIR 基礎建設外銷", 14, DARK_BLUE, True, PP_ALIGN.CENTER),
])

slide_number(slide, 8, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 9: Competitive Landscape — 競爭格局
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "⚔️  競爭格局：FHIR 原生品質量測 = 空白市場", font_size=32, color=WHITE, bold=True)

# 2x2 Matrix
# Axes
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(3), Inches(0.4),
             "← 傳統資料庫                        FHIR 原生 →", font_size=13, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Quadrants
# Top-left: Traditional + Quality
add_shape_bg(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3), RGBColor(0xFF, 0xF0, 0xF0))
add_text_box(slide, Inches(0.7), Inches(2.05), Inches(5.4), Inches(0.3),
             "傳統架構 × 品質分析", font_size=14, color=RED, bold=True)
add_text_box(slide, Inches(0.7), Inches(2.4), Inches(5.4), Inches(1.8),
             "• Health Catalyst（美，市值 8 億美元）— 每院導入 6-12 月\n• Premier Inc.（美，市值 25 億美元）— 4,400 醫院 benchmarking\n• Vizient（美，營收 18 億美元）— 3,000 會員醫院\n• Medisolv（美，小型）— 專做 eCQM 提交\n• Arcadia（美，估值 10 億）— 人口健康管理",
             font_size=13, color=BLACK)

# Top-right: FHIR + Quality → US!
add_shape_bg(slide, Inches(6.5), Inches(2.0), Inches(6.3), Inches(2.3), RGBColor(0xE8, 0xFF, 0xE8))
add_text_box(slide, Inches(6.7), Inches(2.05), Inches(5.9), Inches(0.3),
             "FHIR 原生 × 品質分析", font_size=14, color=ACCENT_GREEN, bold=True)
star = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(2.6), Inches(2.8), Inches(1.4))
star.fill.solid()
star.fill.fore_color.rgb = ACCENT_ORANGE
star.line.fill.background()
add_text_box(slide, Inches(8.2), Inches(2.65), Inches(2.8), Inches(0.6),
             "🔴 我們在這裡", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.2), Inches(3.2), Inches(2.8), Inches(0.6),
             "全球空白市場", font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Bottom-left: Traditional + Platform  
add_shape_bg(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.0), RGBColor(0xF5, 0xF5, 0xFF))
add_text_box(slide, Inches(0.7), Inches(4.55), Inches(5.4), Inches(0.3),
             "傳統架構 × 數據平台", font_size=14, color=GRAY, bold=True)
add_text_box(slide, Inches(0.7), Inches(4.9), Inches(5.4), Inches(1.4),
             "• IQVIA（美，市值 400 億美元）— 處方數據聚合\n• Veeva Systems（美，市值 350 億美元）— 生命科學 CRM\n• Flatiron Health（被 Roche 19 億美元收購）— 腫瘤數據",
             font_size=13, color=BLACK)

# Bottom-right: FHIR + Platform
add_shape_bg(slide, Inches(6.5), Inches(4.5), Inches(6.3), Inches(2.0), RGBColor(0xF0, 0xF8, 0xFF))
add_text_box(slide, Inches(6.7), Inches(4.55), Inches(5.9), Inches(0.3),
             "FHIR 原生 × 基礎設施", font_size=14, color=BRAND_BLUE, bold=True)
add_text_box(slide, Inches(6.7), Inches(4.9), Inches(5.9), Inches(1.4),
             "• Smile CDR（加拿大，估值 3 億美元）— FHIR Server\n• Microsoft Azure FHIR — 雲端基礎設施\n• MITRE CQF — CQL 標準制定（非營利）\n• Lantana Group（美）— FHIR 品質量測顧問",
             font_size=13, color=BLACK)

slide_number(slide, 9, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 10: Why FHIR — 為什麼 FHIR
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BRAND_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "🔑  為什麼 FHIR？為什麼是現在？", font_size=32, color=WHITE, bold=True)

# Three reasons
reasons = [
    ("衛福部強制推行", "次世代電子病歷要求\n所有醫院建 FHIR Server\n\n這不是選配，是法規要求\n\"火車已經開了\"", "🏛️", DARK_BLUE),
    ("TW Core 標準底定", "台灣本土化 FHIR Profile\n已定義完成\n\n我們的 CQL 就是基於\nTW Core 編寫", "📐", BRAND_BLUE),
    ("品質指標電子化", "政府要求品質指標\n要能自動計算\n不再接受手動填報\n\n醫院的剛需", "📊", ACCENT_ORANGE),
]

for i, (title, content, icon, color) in enumerate(reasons):
    x = Inches(0.5 + i * 4.2)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.8), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    add_text_box(slide, x + Inches(0.2), Inches(1.8), Inches(3.4), Inches(0.5),
                 f"{icon}  {title}", font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(2.5),
                 content, font_size=15, color=BLACK)

# Timing quote
add_shape_bg(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8), DARK_BLUE)
add_multi_text(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.6), [
    ("\"兩年前做太早 — 醫院沒有 FHIR Server\"", 18, ACCENT_ORANGE, True, PP_ALIGN.CENTER),
    ("\"兩年後做太晚 — 別人已經卡位\"", 18, ACCENT_ORANGE, True, PP_ALIGN.CENTER),
    ("\"現在是唯一正確的時間點\"", 20, WHITE, True, PP_ALIGN.CENTER),
])

slide_number(slide, 10, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 11: Moat — 護城河
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "🏰  護城河分析", font_size=32, color=WHITE, bold=True)

moats = [
    ("技術門檻  ✅", "CQL 在亞洲幾乎沒人會寫\nFHIR + CQL + 39 指標的\n完整組合是稀缺能力\n\n已申請專利保護"),
    ("網路效應  ⏳", "醫院越多 → 數據越有價值\n→ 更多人買數據\n→ 更多醫院想加入\n\n一旦啟動，飛輪效應極強"),
    ("轉換成本  ✅", "醫院接上系統後\n重新配置 FHIR 映射成本高\n品管人員已習慣操作介面\n\n換掉的痛苦 > 留下的成本"),
    ("數據資產  ⏳", "全台即時品質數據\n一旦累積 100+ 醫院 × 多年\n壟斷地位不可撼動\n\n這是終極護城河"),
]

for i, (title, content) in enumerate(moats):
    x = Inches(0.4 + i * 3.2)
    
    has_it = "✅" in title
    border_color = ACCENT_GREEN if has_it else ACCENT_ORANGE
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.0), Inches(4.0))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = border_color
    card.line.width = Pt(2)
    
    add_text_box(slide, x + Inches(0.2), Inches(1.8), Inches(2.6), Inches(0.5),
                 title, font_size=18, color=border_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.6), Inches(3.0),
                 content, font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)

# Key investor quote
add_shape_bg(slide, Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.3), LIGHT_BLUE)
add_multi_text(slide, Inches(0.8), Inches(6.0), Inches(11.8), Inches(1.1), [
    ("投資人常問：別人能不能抄？", 16, GRAY, False, PP_ALIGN.LEFT),
    ("→ CQL 引擎 + 39 項指標需要深厚臨床品質管理知識，不是純工程問題。已申請專利。而且先行者的數據資產優勢，後進者永遠追不上。", 16, DARK_BLUE, True, PP_ALIGN.LEFT),
])

slide_number(slide, 11, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 12: Milestones — 里程碑
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BRAND_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "🗓️  執行里程碑", font_size=32, color=WHITE, bold=True)

# Timeline
phases = [
    ("Phase 1\n精實驗證", "0-12 個月", 
     "• 約聘模式運營\n• 取得 5-10 家醫院 Pilot\n• 完善產品穩定度\n• 月 burn rate < 15 萬",
     ACCENT_ORANGE, "✅ 已完成產品原型"),
    ("Phase 2\n快速擴張", "12-24 個月",
     "• 擴展至 30-50 家醫院\n• 數據開始有聚合價值\n• 爭取第一筆數據授權營收\n• 評估轉全職時機",
     BRAND_BLUE, "🎯 轉全職觸發點"),
    ("Phase 3\n數據變現", "24-36 個月",
     "• 100+ 家醫院上線\n• 穩定數據授權營收\n• 啟動 Series A 募資\n• 開始評估海外市場",
     ACCENT_GREEN, "💰 Series A Ready"),
    ("Phase 4\n國際擴展", "36+ 個月",
     "• 東南亞市場試點\n• FHIR + 伺服器外銷\n• 地區化 ValueSet 開發\n• 建立亞太區數據聯盟",
     DARK_BLUE, "🌏 亞太佈局"),
]

for i, (title, period, content, color, badge_text) in enumerate(phases):
    x = Inches(0.3 + i * 3.25)
    
    # Timeline line
    if i < 3:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(2.8), Inches(2.3), Inches(0.7), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = GRAY
        line.line.fill.background()
    
    # Circle marker
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), Inches(2.0), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # Period
    add_text_box(slide, x, Inches(1.5), Inches(2.8), Inches(0.4),
                 period, font_size=12, color=GRAY, bold=False, alignment=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide, x, Inches(2.8), Inches(2.8), Inches(0.8),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Content
    add_text_box(slide, x + Inches(0.1), Inches(3.7), Inches(2.6), Inches(2.2),
                 content, font_size=13, color=BLACK)
    
    # Badge
    badge = add_shape_bg(slide, x + Inches(0.2), Inches(5.9), Inches(2.4), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.2), Inches(5.9), Inches(2.4), Inches(0.35),
                 badge_text, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 400萬 runway note
add_shape_bg(slide, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.7), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
             "💡 約聘模式月 burn rate ~15 萬 → 400 萬資金可支撐 26 個月跑道，足夠完成 Phase 1 + 進入 Phase 2",
             font_size=14, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(slide, 12, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 13: Team & Funding — 團隊與融資
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
             "👥  團隊與融資結構", font_size=32, color=WHITE, bold=True)

# Team
team_data = [
    ("CEO / 創辦人", "80%", "FHIR/CQL 技術核心\n醫療品質管理領域專家\n產品與商業策略主導"),
    ("CTO（約聘）", "5%", "系統架構與開發\nFHIR Server 對接\n4 年 vesting + 1 年 cliff"),
    ("CFO（約聘）", "5%", "財務規劃與募資\n醫院合約洽談\n4 年 vesting + 1 年 cliff"),
]

for i, (role, equity, desc) in enumerate(team_data):
    x = Inches(0.5 + i * 4.2)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.8), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BLUE if i == 0 else WHITE
    card.line.color.rgb = BRAND_BLUE
    card.line.width = Pt(1.5)
    
    add_text_box(slide, x + Inches(0.3), Inches(1.8), Inches(3.2), Inches(0.4),
                 role, font_size=18, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.3),
                 f"持股 {equity}", font_size=16, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(1.3),
                 desc, font_size=13, color=BLACK)

# Funding structure
add_shape_bg(slide, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.8), LIGHT_GRAY)
add_multi_text(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.6), [
    ("種子輪融資", 22, DARK_BLUE, True, PP_ALIGN.LEFT),
    ("", 8, BLACK, False, PP_ALIGN.LEFT),
    ("募資金額：NT$ 300 萬", 16, BLACK, False, PP_ALIGN.LEFT),
    ("釋出股權：10%", 16, BLACK, False, PP_ALIGN.LEFT),
    ("Pre-money 估值：NT$ 2,700 萬", 16, BLACK, False, PP_ALIGN.LEFT),
    ("創辦人自有資金：NT$ 100 萬", 16, BLACK, False, PP_ALIGN.LEFT),
    ("總可用資金：NT$ 400 萬", 16, ACCENT_ORANGE, True, PP_ALIGN.LEFT),
])

# Use of funds
add_shape_bg(slide, Inches(6.8), Inches(4.4), Inches(6.0), Inches(2.8), LIGHT_GRAY)
add_multi_text(slide, Inches(7.1), Inches(4.5), Inches(5.5), Inches(2.6), [
    ("資金用途（月 burn rate ~15 萬）", 22, DARK_BLUE, True, PP_ALIGN.LEFT),
    ("", 8, BLACK, False, PP_ALIGN.LEFT),
    ("CTO 約聘：6-8 萬/月", 16, BLACK, False, PP_ALIGN.LEFT),
    ("CFO 約聘：2-3 萬/月", 16, BLACK, False, PP_ALIGN.LEFT),
    ("雲端/伺服器：1-2 萬/月", 16, BLACK, False, PP_ALIGN.LEFT),
    ("法律/專利維護：1 萬/月", 16, BLACK, False, PP_ALIGN.LEFT),
    ("差旅/展覽/雜支：1-2 萬/月", 16, BLACK, False, PP_ALIGN.LEFT),
    ("→ 預估跑道：26 個月", 16, ACCENT_GREEN, True, PP_ALIGN.LEFT),
])

slide_number(slide, 13, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# SLIDE 14: The Ask
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_ORANGE)

add_multi_text(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(6.0), [
    ("The Ask", 44, WHITE, True, PP_ALIGN.CENTER),
    ("", 16, WHITE, False, PP_ALIGN.CENTER),
    ("我們正在尋找種子輪投資 NT$ 300 萬", 28, ACCENT_ORANGE, True, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("用於取得第一批醫院試點合約、建立數據聚合管道", 20, WHITE, False, PP_ALIGN.CENTER),
    ("", 24, WHITE, False, PP_ALIGN.CENTER),
    ("━━━━━━━━━━━━━━━━━━━━━━", 20, GRAY, False, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("我們已經有：", 22, ACCENT_BLUE, True, PP_ALIGN.LEFT),
    ("✅ 完整產品原型 — 39 項品質指標 + 4 大儀表板", 18, WHITE, False, PP_ALIGN.LEFT),
    ("✅ 政府認證 — 衛福部 SMART 50 得獎", 18, WHITE, False, PP_ALIGN.LEFT),
    ("✅ 專利保護 — 已申請專利", 18, WHITE, False, PP_ALIGN.LEFT),
    ("✅ 技術領先 — 台灣唯一 CQL 原生引擎", 18, WHITE, False, PP_ALIGN.LEFT),
    ("✅ 市場時機 — 衛福部強制推 FHIR，醫院必須跟進", 18, WHITE, False, PP_ALIGN.LEFT),
    ("", 16, WHITE, False, PP_ALIGN.CENTER),
    ("我們需要的是第一推力，讓飛輪開始轉動。", 22, ACCENT_ORANGE, True, PP_ALIGN.CENTER),
])

slide_number(slide, 14, TOTAL_SLIDES)


# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
output_path = r"c:\Users\tony1\Desktop\FHIR-CQL-Quality-Platform-20260204\FHIR-CQL_Pitch_Deck.pptx"
prs.save(output_path)
print(f"✅ Pitch deck saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
