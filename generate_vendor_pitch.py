"""
生成「FHIR CQL 平台 × 伺服器廠商合作提案」PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 色彩定義 ──
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xc7)
ACCENT_GREEN = RGBColor(0x10, 0xb9, 0x81)
ACCENT_ORANGE = RGBColor(0xf5, 0x9e, 0x0b)
ACCENT_PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xcc)
DARK_TEXT = RGBColor(0x1e, 0x29, 0x3b)
MID_GRAY = RGBColor(0x64, 0x74, 0x8b)
SOFT_BG = RGBColor(0xf8, 0xfa, 0xfc)
CARD_BG = RGBColor(0xff, 0xff, 0xff)
HIGHLIGHT_BG = RGBColor(0xef, 0xf6, 0xff)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft JhengHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT, line_spacing=1.5, bullet=False, font_name="Microsoft JhengHei"):
    """lines: list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


# ─────────────────────────────────────────────
# SLIDE 1: 封面
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

# 頂部裝飾線
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1),
         "FHIR CQL 醫療品質平台", 44, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
         "軟硬整合合作提案", 32, ACCENT_BLUE, False, PP_ALIGN.CENTER)

# 分隔線
add_rect(slide, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.04), ACCENT_BLUE)

add_text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
         "以 CQL 執行平台 × 伺服器硬體，打造台灣醫療品質基礎設施", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# 底部三個 tag
tags = ["🏆 得獎平台", "🏥 50+ CQL 模組", "🌐 FHIR 4.0.1 標準"]
for i, tag in enumerate(tags):
    x = Inches(3.2 + i * 2.5)
    add_shape(slide, x, Inches(5.0), Inches(2.2), Inches(0.55), RGBColor(0x25, 0x25, 0x45), ACCENT_BLUE, 1)
    add_text(slide, x, Inches(5.05), Inches(2.2), Inches(0.5),
             tag, 14, WHITE, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
         "2026 年 4 月", 14, MID_GRAY, False, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 2: 市場痛點
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "🔍 市場痛點：CQL 寫好了，卻無處可跑", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_BLUE)

# 痛點卡片
pain_points = [
    ("😰", "CQL 創作者的困境",
     "衛福部推動 FHIR + CQL\n各單位撰寫品質指標 CQL\n但沒有平台可以實際執行"),
    ("🔒", "醫院的資安顧慮",
     "醫院不可能開放 FHIR Server\n給每個 CQL 創作者遠端撈資料\n資料外洩風險太大"),
    ("📊", "品質量測的現況",
     "人工抽審 + Excel 彙整\n耗時、不即時、難以標準化\n400+ 家醫院都在用土法煉鋼"),
]
for i, (icon, title, desc) in enumerate(pain_points):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.6), Inches(3.7), Inches(3.8), CARD_BG, RGBColor(0xe2, 0xe8, 0xf0), 1)
    add_text(slide, x + Inches(0.3), Inches(1.8), Inches(3.1), Inches(0.7), icon, 40, DARK_TEXT, False, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(2.6), Inches(3.1), Inches(0.5), title, 20, DARK_TEXT, True, PP_ALIGN.CENTER)
    add_multiline(slide, x + Inches(0.3), Inches(3.3), Inches(3.1), Inches(2.0),
                  desc.split('\n'), 15, MID_GRAY, 1.8)

# 結論
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), HIGHLIGHT_BG, ACCENT_BLUE, 1.5)
add_text(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.9),
         "💡 結論：台灣需要一個「CQL 執行基礎設施」— 讓指標在醫院內安全運行，資料不出院", 18, ACCENT_BLUE, True, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 3: 我們的解決方案
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "💡 解決方案：CQL Marketplace + 院內執行引擎", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

# 架構圖 - 用方塊模擬
# 左邊：CQL Marketplace
add_shape(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(5.0), RGBColor(0xec, 0xfd, 0xf5), ACCENT_GREEN, 2)
add_text(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.5),
         "☁️  CQL Marketplace（雲端）", 16, ACCENT_GREEN, True, PP_ALIGN.CENTER)

marketplace_items = [
    "📋 品質指標 CQL 上架",
    "✅ CQL 審核認證機制",
    "👨‍⚕️ 創作者社群",
    "📦 版本管理 & 更新推送",
    "💰 授權計費系統",
]
for j, item in enumerate(marketplace_items):
    add_shape(slide, Inches(1.2), Inches(2.4 + j * 0.75), Inches(2.7), Inches(0.55), CARD_BG, RGBColor(0xd1, 0xfa, 0xe5), 1)
    add_text(slide, Inches(1.3), Inches(2.42 + j * 0.75), Inches(2.5), Inches(0.5), item, 13, DARK_TEXT, False)

# 中間箭頭
add_text(slide, Inches(4.5), Inches(3.5), Inches(1), Inches(1),
         "⟹", 48, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(4.3), Inches(4.3), Inches(1.4), Inches(0.5),
         "CQL 下載\n結果回傳", 11, MID_GRAY, False, PP_ALIGN.CENTER)

# 右邊：院內執行環境
add_shape(slide, Inches(5.8), Inches(1.6), Inches(6.8), Inches(5.0), HIGHLIGHT_BG, ACCENT_BLUE, 2)
add_text(slide, Inches(5.8), Inches(1.7), Inches(6.8), Inches(0.5),
         "🏥  醫院內部（伺服器 Appliance）", 16, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# 伺服器內部元件
components = [
    (Inches(6.2), Inches(2.5), "🖥️ CQL 執行引擎", CARD_BG, Inches(2.8)),
    (Inches(9.3), Inches(2.5), "🔗 FHIR 資料適配器", CARD_BG, Inches(2.8)),
    (Inches(6.2), Inches(3.6), "📊 品質報表產生器", CARD_BG, Inches(2.8)),
    (Inches(9.3), Inches(3.6), "🔐 資安認證模組", CARD_BG, Inches(2.8)),
]
for (x, y, label, bg, w) in components:
    add_shape(slide, x, y, w, Inches(0.7), bg, ACCENT_BLUE, 1)
    add_text(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.5), label, 14, DARK_TEXT, False, PP_ALIGN.CENTER)

# 資料來源
add_text(slide, Inches(6.2), Inches(4.7), Inches(6), Inches(0.4),
         "資料來源（不出院）：", 13, MID_GRAY, True)
sources = ["FHIR Server", "HIS 資料庫", "CSV 匯入"]
for i, src in enumerate(sources):
    add_shape(slide, Inches(6.2 + i * 2.1), Inches(5.2), Inches(1.9), Inches(0.55), RGBColor(0xe0, 0xf2, 0xfe), ACCENT_BLUE, 1)
    add_text(slide, Inches(6.2 + i * 2.1), Inches(5.22), Inches(1.9), Inches(0.5), src, 13, ACCENT_BLUE, False, PP_ALIGN.CENTER)

# 底部重點
add_text(slide, Inches(6.2), Inches(6.0), Inches(6), Inches(0.5),
         "🔒 資料永遠留在醫院內部，只有統計結果對外", 14, ACCENT_GREEN, True)


# ─────────────────────────────────────────────
# SLIDE 4: 合作模式
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_PURPLE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "🤝 合作模式：軟硬整合，各司其職", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_PURPLE)

# 左邊：我方
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5), RGBColor(0xf0, 0xf0, 0xff), ACCENT_PURPLE, 2)
add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6),
         "🧠 我方：平台 & CQL 生態", 20, ACCENT_PURPLE, True, PP_ALIGN.CENTER)

my_items = [
    ("CQL 執行引擎", "核心技術，50+ 品質指標模組"),
    ("Marketplace 平台", "CQL 上架、審核、版本管理"),
    ("FHIR 資料適配器", "支援 FHIR / SQL / CSV 多種來源"),
    ("品質報表 & 儀表板", "即時視覺化分析介面"),
    ("技術支援 & 更新", "持續維護、指標更新推送"),
]
for j, (title, desc) in enumerate(my_items):
    y = Inches(2.3 + j * 0.7)
    add_text(slide, Inches(1.2), y, Inches(4.8), Inches(0.35), f"✦  {title}", 15, DARK_TEXT, True)
    add_text(slide, Inches(1.7), y + Inches(0.3), Inches(4.3), Inches(0.3), desc, 12, MID_GRAY)

# 右邊：貴方
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.5), RGBColor(0xf0, 0xf8, 0xff), ACCENT_BLUE, 2)
add_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.6),
         "🖥️ 貴方：伺服器 & 通路", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER)

their_items = [
    ("伺服器硬體", "醫療等級伺服器、穩定運行環境"),
    ("醫院通路 & 業務團隊", "現有客戶關係、銷售網絡"),
    ("機房 & 維運支援", "硬體安裝、網路配置、SLA 保障"),
    ("資安合規", "符合醫院 IT 規範、ISO 27001"),
    ("售後服務", "硬體保固、現場技術支援"),
]
for j, (title, desc) in enumerate(their_items):
    y = Inches(2.3 + j * 0.7)
    add_text(slide, Inches(7.4), y, Inches(4.8), Inches(0.35), f"✦  {title}", 15, DARK_TEXT, True)
    add_text(slide, Inches(7.9), y + Inches(0.3), Inches(4.3), Inches(0.3), desc, 12, MID_GRAY)

# 底部整合
add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), ACCENT_PURPLE)
add_text(slide, Inches(1.2), Inches(6.3), Inches(11), Inches(0.7),
         "🎯 整合成果：醫院買一台伺服器 → 開箱即有 CQL 品質量測平台 → 連上 Marketplace 下載指標 → 即刻運行", 16, WHITE, True, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 5: 商業模式 & 分潤
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_ORANGE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "💰 商業模式 & 營收結構", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_ORANGE)

# 營收來源表格
headers = ["營收來源", "付費方", "模式", "分配"]
rows = [
    ["伺服器硬體銷售", "醫院", "一次性購買", "貴方為主"],
    ["平台授權費（年費）", "醫院", "年度訂閱", "我方為主"],
    ["CQL 模組費", "醫院", "按模組訂閱", "我方抽成 + 創作者"],
    ["客製化開發", "醫院", "專案計費", "雙方合作"],
    ["維護保固", "醫院", "年度合約", "各自負責範圍"],
]

# 表頭
for i, h in enumerate(headers):
    x = Inches(0.8 + i * 3.05)
    add_shape(slide, x, Inches(1.5), Inches(2.85), Inches(0.6), ACCENT_ORANGE)
    add_text(slide, x, Inches(1.55), Inches(2.85), Inches(0.5), h, 15, WHITE, True, PP_ALIGN.CENTER)

# 表格內容
for r, row in enumerate(rows):
    bg = CARD_BG if r % 2 == 0 else RGBColor(0xfe, 0xf9, 0xef)
    for c, cell in enumerate(row):
        x = Inches(0.8 + c * 3.05)
        y = Inches(2.1 + r * 0.6)
        add_shape(slide, x, y, Inches(2.85), Inches(0.55), bg, RGBColor(0xe2, 0xe8, 0xf0), 0.5)
        add_text(slide, x, y + Inches(0.05), Inches(2.85), Inches(0.45), cell, 13, DARK_TEXT, False, PP_ALIGN.CENTER)

# 市場規模
add_text(slide, Inches(0.8), Inches(5.2), Inches(12), Inches(0.5),
         "📈 目標市場規模", 20, DARK_TEXT, True)

markets = [
    ("醫學中心", "~20 家", "已有 FHIR", ACCENT_BLUE),
    ("區域醫院", "~80 家", "部分有 FHIR", ACCENT_GREEN),
    ("地區醫院", "~300 家", "CSV/SQL 模式", ACCENT_ORANGE),
    ("合計潛力", "400+ 家", "全台覆蓋", ACCENT_PURPLE),
]
for i, (seg, num, note, clr) in enumerate(markets):
    x = Inches(0.8 + i * 3.1)
    add_shape(slide, x, Inches(5.7), Inches(2.8), Inches(1.2), CARD_BG, clr, 2)
    add_text(slide, x, Inches(5.8), Inches(2.8), Inches(0.4), seg, 15, clr, True, PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.15), Inches(2.8), Inches(0.4), num, 22, DARK_TEXT, True, PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.5), Inches(2.8), Inches(0.3), note, 11, MID_GRAY, False, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 6: 合作時程
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "📅 合作推進時程", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_BLUE)

# 時間軸
phases = [
    ("Phase 1", "POC 驗證", "1-2 個月",
     ["選定 1-2 家合作醫院", "整合伺服器 + 平台 Demo 機", "實際跑 3-5 個品質指標", "收集醫院回饋"],
     ACCENT_BLUE),
    ("Phase 2", "產品化", "3-4 個月",
     ["完善 Appliance 安裝流程", "建立 CQL Marketplace MVP", "制定定價 & 分潤合約", "取得資安認證"],
     ACCENT_GREEN),
    ("Phase 3", "市場推廣", "5-12 個月",
     ["醫學中心率先導入", "參加 HIMSS / 醫療資訊展", "擴展到區域醫院", "建立 CQL 創作者社群"],
     ACCENT_ORANGE),
]

for i, (phase, title, duration, items, clr) in enumerate(phases):
    x = Inches(0.8 + i * 4.2)
    # 卡片
    add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(5.2), CARD_BG, clr, 2)
    # Phase 標籤
    add_shape(slide, x + Inches(0.2), Inches(1.7), Inches(1.5), Inches(0.45), clr)
    add_text(slide, x + Inches(0.2), Inches(1.72), Inches(1.5), Inches(0.4), phase, 14, WHITE, True, PP_ALIGN.CENTER)
    # 時間
    add_text(slide, x + Inches(1.9), Inches(1.72), Inches(1.6), Inches(0.4), duration, 13, MID_GRAY, False, PP_ALIGN.RIGHT)
    # 標題
    add_text(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.5), title, 22, DARK_TEXT, True)
    # 項目
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.3), Inches(3.2 + j * 0.65), Inches(3.2), Inches(0.5),
                 f"▸ {item}", 14, MID_GRAY)

# 底部
add_shape(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.0), SOFT_BG)


# ─────────────────────────────────────────────
# SLIDE 7: 我方優勢 & 得獎背書
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "🏆 為什麼是我們？", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_GREEN)

advantages = [
    ("🏆", "得獎肯定", "產品已獲獎項認可\n具備市場公信力"),
    ("📊", "50+ CQL 模組", "涵蓋疾管、國健、ESG、品質\n四大類別即可使用"),
    ("🌐", "國際標準", "基於 FHIR 4.0.1 + CQL\n完全符合衛福部推動方向"),
    ("🔌", "多源適配", "不需等醫院建 FHIR Server\nCSV / SQL / FHIR 都能接"),
    ("🚀", "即戰力", "平台已可運行展示\n不是 PPT 階段，是產品階段"),
    ("🎯", "零競爭者", "台灣目前沒有公開的\nCQL 執行平台，先行者優勢"),
]

for i, (icon, title, desc) in enumerate(advantages):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.5 + row * 2.8)
    add_shape(slide, x, y, Inches(3.7), Inches(2.4), CARD_BG, ACCENT_GREEN, 1.5)
    add_text(slide, x, y + Inches(0.2), Inches(3.7), Inches(0.6), icon, 36, DARK_TEXT, False, PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.85), Inches(3.7), Inches(0.4), title, 18, DARK_TEXT, True, PP_ALIGN.CENTER)
    add_multiline(slide, x + Inches(0.3), y + Inches(1.35), Inches(3.1), Inches(0.9),
                  desc.split('\n'), 13, MID_GRAY, 1.6)


# ─────────────────────────────────────────────
# SLIDE 8: 市場分析
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
ACCENT_RED = RGBColor(0xef, 0x44, 0x44)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_RED)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "📈 市場分析：數據說話", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_RED)

# ── 台灣醫療市場大數字 ──
add_text(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.5),
         "台灣醫療市場規模（衛福部統計）", 18, DARK_TEXT, True)

big_nums = [
    ("8,700 億+", "健保年度支出\n（2025年）", ACCENT_RED),
    ("475 家", "全台醫院\n（醫學中心+區域+地區）", ACCENT_BLUE),
    ("23,000+", "基層診所\n（未來擴展市場）", ACCENT_GREEN),
    ("2,400 萬", "健保涵蓋人口\n（覆蓋率 99.9%）", ACCENT_PURPLE),
]
for i, (num, label, clr) in enumerate(big_nums):
    x = Inches(0.8 + i * 3.1)
    add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(1.5), CARD_BG, clr, 2)
    add_text(slide, x, Inches(1.9), Inches(2.8), Inches(0.6), num, 26, clr, True, PP_ALIGN.CENTER)
    add_multiline(slide, x + Inches(0.2), Inches(2.55), Inches(2.4), Inches(0.7),
                  label.split('\n'), 12, MID_GRAY, 1.4)

# ── 目標客群分層 ──
add_text(slide, Inches(0.8), Inches(3.6), Inches(6), Inches(0.5),
         "目標客群與預估營收（保守估計）", 18, DARK_TEXT, True)

# 表頭
seg_headers = ["客群", "家數", "年費/家", "滲透率", "年營收估計"]
for i, h in enumerate(seg_headers):
    widths = [Inches(2.2), Inches(1.5), Inches(2.2), Inches(1.5), Inches(2.6)]
    x_pos = Inches(0.8)
    for k in range(i):
        x_pos += widths[k]
    add_shape(slide, x_pos, Inches(4.1), widths[i], Inches(0.5), ACCENT_RED)
    add_text(slide, x_pos, Inches(4.13), widths[i], Inches(0.45), h, 13, WHITE, True, PP_ALIGN.CENTER)

# 表格資料
seg_rows = [
    ["醫學中心", "19 家", "50-80 萬/年", "30%（~6家）", "300-480 萬"],
    ["區域醫院", "86 家", "30-50 萬/年", "15%（~13家）", "390-650 萬"],
    ["地區醫院", "370 家", "15-30 萬/年", "10%（~37家）", "555-1,110 萬"],
    ["合計（3年目標）", "475 家", "—", "~56 家", "1,245-2,240 萬/年"],
]
for r, row in enumerate(seg_rows):
    is_total = (r == len(seg_rows) - 1)
    bg = RGBColor(0xff, 0xf1, 0xf2) if is_total else (CARD_BG if r % 2 == 0 else RGBColor(0xfe, 0xf9, 0xef))
    x_pos = Inches(0.8)
    for c, cell in enumerate(row):
        widths = [Inches(2.2), Inches(1.5), Inches(2.2), Inches(1.5), Inches(2.6)]
        add_shape(slide, x_pos, Inches(4.6 + r * 0.5), widths[c], Inches(0.48), bg, RGBColor(0xe2, 0xe8, 0xf0), 0.5)
        clr_cell = ACCENT_RED if (is_total and c == 4) else DARK_TEXT
        bld = is_total
        add_text(slide, x_pos, Inches(4.63 + r * 0.5), widths[c], Inches(0.42), cell, 12, clr_cell, bld, PP_ALIGN.CENTER)
        x_pos += widths[c]

# ── 政策驅動力 ──
add_text(slide, Inches(0.8), Inches(6.7), Inches(12), Inches(0.5),
         "🔥 政策風口：衛福部 2024 起要求醫學中心導入 FHIR｜2026 擴大至區域醫院｜品質指標申報將數位化", 15, ACCENT_RED, True)


# ─────────────────────────────────────────────
# SLIDE 9: 競爭者分析
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xf1)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_INDIGO)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "⚔️ 競爭者分析：台灣 CQL 執行平台 = 零", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_INDIGO)

# ── 競爭對手表格 ──
comp_headers = ["競爭者 / 替代方案", "類型", "CQL\n執行", "FHIR\n整合", "台灣\n在地化", "品質指標\n自動量測", "威脅程度"]
comp_widths = [Inches(2.8), Inches(1.4), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.2), Inches(1.3)]

for i, h in enumerate(comp_headers):
    x_pos = Inches(0.8)
    for k in range(i):
        x_pos += comp_widths[k]
    add_shape(slide, x_pos, Inches(1.4), comp_widths[i], Inches(0.7), ACCENT_INDIGO)
    add_multiline(slide, x_pos + Inches(0.05), Inches(1.42), comp_widths[i] - Inches(0.1), Inches(0.65),
                  h.split('\n'), 11, WHITE, 1.2)

# 競爭者資料
competitors = [
    ["台灣 HIS 廠商\n（衛生、昕力、醫揚）", "HIS 系統商", "❌", "⚠️ 起步", "✅", "❌ 人工", "中"],
    ["關貿網路\n（EMR 交換平台）", "政府承包商", "❌", "⚠️ 部分", "✅", "❌", "低"],
    ["Smile Digital Health\n（加拿大，CQF-Ruler）", "國際 FHIR 廠", "✅", "✅", "❌", "⚠️ 通用", "低"],
    ["Epic / Oracle Health\n（美國 EHR 巨頭）", "國際 EHR", "✅ 內建", "✅", "❌", "✅ 美國", "極低"],
    ["醫院自行開發\n（Excel + 人工統計）", "現行做法", "❌", "❌", "✅", "❌ 人工", "替代品"],
    ["我們（FHIR CQL 平台）", "專注平台", "✅", "✅", "✅", "✅ 50+", "—"],
]
comp_colors_threat = {"極低": MID_GRAY, "低": ACCENT_GREEN, "中": ACCENT_ORANGE, "替代品": ACCENT_RED, "—": ACCENT_INDIGO}

for r, row in enumerate(competitors):
    is_us = (r == len(competitors) - 1)
    bg = RGBColor(0xe8, 0xe8, 0xff) if is_us else (CARD_BG if r % 2 == 0 else RGBColor(0xf8, 0xfa, 0xfc))
    border = ACCENT_INDIGO if is_us else RGBColor(0xe2, 0xe8, 0xf0)
    bw = 2 if is_us else 0.5
    x_pos = Inches(0.8)
    for c, cell in enumerate(row):
        add_shape(slide, x_pos, Inches(2.1 + r * 0.65), comp_widths[c], Inches(0.63), bg, border, bw)
        # 多行文字
        cell_color = ACCENT_INDIGO if is_us else DARK_TEXT
        if c == 6:  # 威脅程度欄位上色
            cell_color = comp_colors_threat.get(cell, DARK_TEXT)
        cell_bold = is_us or c == 6
        add_multiline(slide, x_pos + Inches(0.08), Inches(2.12 + r * 0.65),
                      comp_widths[c] - Inches(0.16), Inches(0.58),
                      cell.split('\n'), 11, cell_color, 1.1)
        x_pos += comp_widths[c]

# ── 結論 ──
add_shape(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.1), CARD_BG, ACCENT_INDIGO, 2)

conclusion_lines = [
    ("關鍵發現", True, ACCENT_INDIGO),
    ("① 台灣目前沒有任何公開的 CQL 執行平台 — 我們是唯一", False, DARK_TEXT),
    ("② 國際廠商有技術但不了解台灣法規、編碼、指標定義 — 在地化是壁壘", False, DARK_TEXT),
    ("③ 現行替代方案（Excel 人工）效率極低 — 痛點已被驗證", False, DARK_TEXT),
]
add_multiline(slide, Inches(1.2), Inches(6.15), Inches(10.8), Inches(1.0),
              conclusion_lines, 14, DARK_TEXT, 1.5)


# ─────────────────────────────────────────────
# SLIDE 10: 未來擴展 — CDS Hooks 即時決策支援
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SOFT_BG)
ACCENT_TEAL = RGBColor(0x06, 0xb6, 0xd4)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_TEAL)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "🔮 未來擴展：從品質量測到即時臨床決策", 30, DARK_TEXT, True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_TEAL)

# 左：現在（品質量測）
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.2), RGBColor(0xec, 0xfd, 0xf5), ACCENT_GREEN, 2)
add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
         "📊 現在：品質量測引擎（eCQM）", 18, ACCENT_GREEN, True, PP_ALIGN.CENTER)

now_items = [
    ("批次分析", "撈一整季資料計算指標"),
    ("群體層級", "全院抗生素使用率、剖腹產率"),
    ("報表輸出", "統計報表 + 可追溯稽核"),
    ("管理者使用", "品質管理部門、衛福部申報"),
]
for j, (title, desc) in enumerate(now_items):
    y = Inches(2.3 + j * 0.6)
    add_text(slide, Inches(1.2), y, Inches(4.8), Inches(0.3), f"✅  {title}", 14, DARK_TEXT, True)
    add_text(slide, Inches(2.2), y + Inches(0.02), Inches(3.8), Inches(0.3), f"— {desc}", 13, MID_GRAY)

# 中間箭頭
add_text(slide, Inches(6.5), Inches(2.6), Inches(0.8), Inches(0.8),
         "＋", 48, ACCENT_TEAL, True, PP_ALIGN.CENTER)

# 右：未來（CDS Hooks）
add_shape(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(3.2), HIGHLIGHT_BG, ACCENT_TEAL, 2)
add_text(slide, Inches(7.3), Inches(1.6), Inches(5.5), Inches(0.5),
         "⚡ 未來：CDS Hooks 即時決策（CDS）", 18, ACCENT_TEAL, True, PP_ALIGN.CENTER)

future_items = [
    ("即時觸發", "醫生開藥時自動提醒"),
    ("個人層級", "這個病人現在該不該用這個藥？"),
    ("卡片推送", "建議卡片直接顯示在 EHR 畫面"),
    ("臨床醫師使用", "即時輔助臨床決策"),
]
for j, (title, desc) in enumerate(future_items):
    y = Inches(2.3 + j * 0.6)
    add_text(slide, Inches(7.7), y, Inches(4.8), Inches(0.3), f"🔜  {title}", 14, DARK_TEXT, True)
    add_text(slide, Inches(8.7), y + Inches(0.02), Inches(3.8), Inches(0.3), f"— {desc}", 13, MID_GRAY)

# CDS Hooks 流程圖（用方塊模擬）
add_text(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(0.5),
         "⚙️ CDS Hooks 運作流程", 18, DARK_TEXT, True)

flow_steps = [
    ("👨‍⚕️ 醫生開藥", RGBColor(0xe0, 0xf2, 0xfe)),
    ("🔗 EHR 觸發 Hook", RGBColor(0xfe, 0xf9, 0xef)),
    ("🧠 CQL 引擎判斷", RGBColor(0xec, 0xfd, 0xf5)),
    ("📋 回傳建議卡片", RGBColor(0xf0, 0xf0, 0xff)),
    ("✅ 醫生採納/忽略", RGBColor(0xff, 0xf1, 0xf2)),
]
for i, (label, bg) in enumerate(flow_steps):
    x = Inches(0.8 + i * 2.5)
    add_shape(slide, x, Inches(5.5), Inches(2.2), Inches(0.7), bg, ACCENT_TEAL, 1)
    add_text(slide, x, Inches(5.53), Inches(2.2), Inches(0.6), label, 13, DARK_TEXT, False, PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_text(slide, x + Inches(2.2), Inches(5.5), Inches(0.3), Inches(0.7),
                 "→", 20, ACCENT_TEAL, True, PP_ALIGN.CENTER)

# 底部重點
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), ACCENT_TEAL)
add_text(slide, Inches(1.2), Inches(6.55), Inches(11), Inches(0.6),
         "🎯 同一個 CQL Marketplace，兩條產品線：品質量測 CQL（現在） ＋ 臨床決策 CQL（未來）", 16, WHITE, True, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 9: 下一步 & 聯繫
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.8),
         "🚀 下一步", 36, WHITE, True, PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(1.9), Inches(2.3), Inches(0.04), ACCENT_BLUE)

next_steps = [
    "1️⃣  安排技術對接會議，評估整合可行性",
    "2️⃣  選定 POC 醫院，組建 Demo 機",
    "3️⃣  簽訂合作意向書（MOU）",
    "4️⃣  啟動 Phase 1 — 用成果說話",
]
for j, step in enumerate(next_steps):
    add_shape(slide, Inches(3), Inches(2.5 + j * 0.85), Inches(7.3), Inches(0.65),
              RGBColor(0x25, 0x25, 0x45), ACCENT_BLUE, 1)
    add_text(slide, Inches(3.3), Inches(2.55 + j * 0.85), Inches(6.7), Inches(0.55),
             step, 17, WHITE, False)

# 願景
add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.6),
         "「每家醫院一台，裡面跑我們的平台」", 22, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
         "讓我們一起打造台灣醫療品質的基礎設施", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ── 儲存 ──
output_path = os.path.join(os.path.dirname(__file__), "FHIR_CQL_合作提案_V2.pptx")
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
