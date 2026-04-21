# -*- coding: utf-8 -*-
"""
合作提案簡報 - 給 SAS 團隊的合作方案 PPT
FHIR-CQL 品質平台 × SAS 數據團隊 合作提案
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 色彩定義 ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x0F, 0x2B, 0x46)
BRAND_BLUE = RGBColor(0x00, 0x6D, 0xAA)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x00, 0xA6, 0x5A)
ACCENT_RED = RGBColor(0xE8, 0x3E, 0x3E)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_PURPLE = RGBColor(0x7B, 0x43, 0x9A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helper Functions ──

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft JhengHei'
    p.alignment = align
    return txBox


def add_multi_text(slide, left, top, width, height, lines):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft JhengHei'
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_table(slide, left, top, width, height, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Microsoft JhengHei'
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRAND_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


# ========================================================
# Slide 1: 封面
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
# 上方裝飾線
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2),
         "QAlert", 52, ACCENT_BLUE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(0.8),
         "即時醫療品質警示系統", 36, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
         "Real-time Clinical Quality Alert System", 22, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
         "FHIR-CQL 品質平台 × SAS 數據團隊 合作提案", 24, ACCENT_BLUE, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
         "從數據到決策的完整解決方案", 20, GRAY, False, PP_ALIGN.CENTER)

# 底部裝飾
add_shape(slide, Inches(0), Inches(6.8), prs.slide_width, Inches(0.7), RGBColor(0x08, 0x1C, 0x30))
add_text(slide, Inches(1.5), Inches(6.85), Inches(10), Inches(0.5),
         "2026 年 4 月", 16, GRAY, False, PP_ALIGN.CENTER)


# ========================================================
# Slide 2: 現況 - 架構圖（對方已熟悉的那張圖）
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "現況架構：你們的左線 × 我們的右線", 30, DARK_BLUE, True)

# 左線框
add_rounded_shape(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(5.3), LIGHT_BLUE, BRAND_BLUE)
add_text(slide, Inches(1.0), Inches(1.6), Inches(4.8), Inches(0.5),
         "左線（貴團隊）", 22, BRAND_BLUE, True, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(1.2), Inches(2.3), Inches(4.5), Inches(4.2), [
    ("HIS 資料", 18, DARK_TEXT, True, PP_ALIGN.CENTER),
    ("↓", 20, GRAY, False, PP_ALIGN.CENTER),
    ("數據中台", 18, DARK_TEXT, True, PP_ALIGN.CENTER),
    ("↓", 20, GRAY, False, PP_ALIGN.CENTER),
    ("SAS EG（整理好的資料庫）", 18, DARK_TEXT, True, PP_ALIGN.CENTER),
    ("↓", 20, GRAY, False, PP_ALIGN.CENTER),
    ("SAS VA 院內儀表板", 18, BRAND_BLUE, True, PP_ALIGN.CENTER),
    ("", 8, GRAY, False, PP_ALIGN.CENTER),
    ("✅ 資料清洗完整", 14, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("✅ 院內視覺化成熟", 14, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("❌ 無法跨院", 14, ACCENT_RED, False, PP_ALIGN.LEFT),
    ("❌ 無法即時警示", 14, ACCENT_RED, False, PP_ALIGN.LEFT),
])

# 右線框
add_rounded_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xF0, 0xFF, 0xF0), ACCENT_GREEN)
add_text(slide, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.5),
         "右線（我們）", 22, ACCENT_GREEN, True, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(4.2), [
    ("FHIR（TW Core 標準格式）", 18, DARK_TEXT, True, PP_ALIGN.CENTER),
    ("↓", 20, GRAY, False, PP_ALIGN.CENTER),
    ("CQL 品質規則引擎", 18, DARK_TEXT, True, PP_ALIGN.CENTER),
    ("↓", 20, GRAY, False, PP_ALIGN.CENTER),
    ("品質儀表板 + 即時警示 App", 18, ACCENT_GREEN, True, PP_ALIGN.CENTER),
    ("↓", 20, GRAY, False, PP_ALIGN.CENTER),
    ("院際交換 / CDS 臨床決策", 18, ACCENT_GREEN, True, PP_ALIGN.CENTER),
    ("", 8, GRAY, False, PP_ALIGN.CENTER),
    ("✅ 國際標準，可跨院", 14, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("✅ 同一套 CQL 可報表 + 即時", 14, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("✅ 已有專利平台", 14, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("⚠️ 需要乾淨的資料來源", 14, ACCENT_ORANGE, False, PP_ALIGN.LEFT),
])

# 中間連接箭頭
add_text(slide, Inches(5.6), Inches(3.8), Inches(1.6), Inches(0.8),
         "合作\n互補", 20, ACCENT_ORANGE, True, PP_ALIGN.CENTER)


# ========================================================
# Slide 3: 合作價值 - 1+1 > 2
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "合作價值：1 + 1 > 2", 30, DARK_BLUE, True)

# 三個價值卡片
cards = [
    ("🔍 交叉驗證", "SAS 程式 vs CQL 規則\n逐指標比對結果",
     "雙方都能發現自己的 bug\n共同提升品質指標準確度", BRAND_BLUE),
    ("🔗 資料串接", "SAS EG 整理好的資料\n轉成 FHIR 國際標準",
     "你們負責清洗（最苦的活）\n我們負責標準化（最難的活）", ACCENT_GREEN),
    ("🚀 功能互補", "SAS VA 看過去\nCQL App 管現在",
     "報表 + 即時警示 = 完整解決方案\n從事後統計到事前預警", ACCENT_ORANGE),
]

for i, (title, desc, value, color) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    add_rounded_shape(slide, x, Inches(1.5), Inches(3.8), Inches(5.0), LIGHT_GRAY, color)
    add_text(slide, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.6),
             title, 22, color, True, PP_ALIGN.CENTER)
    add_shape(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.03), color)
    add_multi_text(slide, x + Inches(0.3), Inches(2.7), Inches(3.2), Inches(1.5), [
        (line, 15, DARK_TEXT, False, PP_ALIGN.LEFT) for line in desc.split('\n')
    ])
    add_multi_text(slide, x + Inches(0.3), Inches(4.2), Inches(3.2), Inches(2.0), [
        (line, 14, GRAY, False, PP_ALIGN.LEFT) for line in value.split('\n')
    ])


# ========================================================
# Slide 4: 合作分工表
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "合作分工：各自發揮專長", 30, DARK_BLUE, True)

data = [
    ["工作項目", "貴團隊（SAS）", "我們（FHIR-CQL）", "產出"],
    ["HIS 資料擷取", "✅ 負責", "—", "原始資料"],
    ["資料清洗整理", "✅ SAS EG", "—", "乾淨的表格"],
    ["轉換成 FHIR 格式", "提供欄位清單", "✅ 寫轉換腳本", "FHIR Resources"],
    ["品質指標規則", "提供 SAS 程式邏輯", "✅ 撰寫 CQL", "標準化規則"],
    ["交叉驗證", "✅ SAS 計算結果", "✅ CQL 計算結果", "驗證報告"],
    ["院內儀表板", "✅ SAS VA", "✅ 品質平台", "雙軌展示"],
    ["即時警示 App", "—", "✅ 開發", "開診前提醒"],
    ["跨院資料交換", "—", "✅ FHIR 標準", "院際互通"],
]

add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), data,
          col_widths=[Inches(2.5), Inches(3.0), Inches(3.5), Inches(3.3)])


# ========================================================
# Slide 5: 即時警示 - 同一套 CQL 兩種用途
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "核心亮點：同一條 CQL 規則，兩種用途", 30, DARK_BLUE, True)

# 中間 CQL box
add_rounded_shape(slide, Inches(4.5), Inches(1.3), Inches(4.3), Inches(1.0), BRAND_BLUE)
add_text(slide, Inches(4.7), Inches(1.4), Inches(3.9), Inches(0.8),
         "同一條 CQL 規則\n（例：糖尿病 HbA1c 檢測）", 16, WHITE, True, PP_ALIGN.CENTER)

# 左邊：事後報表
add_rounded_shape(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(4.0), LIGHT_BLUE, BRAND_BLUE)
add_text(slide, Inches(1.0), Inches(2.9), Inches(5.1), Inches(0.5),
         "用途 1：品質報表（事後統計）", 20, BRAND_BLUE, True, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(1.2), Inches(3.6), Inches(5.0), Inches(3.0), [
    ("▶ 定時排程 / 每月每季", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("▶ 撈全院所有病人資料", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("▶ 計算比率、趨勢圖表", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("", 8, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("產出：「上季 HbA1c 達標率 72%」", 15, BRAND_BLUE, True, PP_ALIGN.LEFT),
    ("使用者：品管師、院長", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 8, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("💡 跟 SAS VA 功能類似", 14, ACCENT_ORANGE, False, PP_ALIGN.LEFT),
    ("　 但規則是國際標準、可跨院", 14, ACCENT_ORANGE, False, PP_ALIGN.LEFT),
])

# 右邊：即時警示
add_rounded_shape(slide, Inches(7.0), Inches(2.8), Inches(5.5), Inches(4.0), RGBColor(0xFF, 0xF3, 0xE0), ACCENT_ORANGE)
add_text(slide, Inches(7.2), Inches(2.9), Inches(5.1), Inches(0.5),
         "用途 2：即時警示（事前預警）", 20, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(7.4), Inches(3.6), Inches(5.0), Inches(3.0), [
    ("▶ 每天凌晨跑 / 開診前完成", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("▶ 只撈今日門診病人", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("▶ 產出個人化警示清單", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("", 8, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("產出：「此病人 HbA1c 逾期 243 天」", 15, ACCENT_ORANGE, True, PP_ALIGN.LEFT),
    ("使用者：第一線醫師、護理師", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 8, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("🔥 SAS VA 做不到！", 14, ACCENT_RED, True, PP_ALIGN.LEFT),
    ("　 因為 SAS 不是即時系統", 14, ACCENT_RED, False, PP_ALIGN.LEFT),
])


# ========================================================
# Slide 6: 即時警示 App 畫面示意
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "即時警示 App：護理站第二個視窗", 30, DARK_BLUE, True)

# 模擬警示畫面
add_rounded_shape(slide, Inches(1.5), Inches(1.3), Inches(10.3), Inches(5.8), WHITE, BRAND_BLUE)

# Header
add_shape(slide, Inches(1.5), Inches(1.3), Inches(10.3), Inches(0.8), BRAND_BLUE)
add_text(slide, Inches(1.8), Inches(1.35), Inches(9.5), Inches(0.7),
         "📋  今日門診品質提醒　　2026-04-14 上午診　　內分泌科 王醫師", 18, WHITE, True)

# 警示列表
alerts = [
    ("09:00", "王大明", "🔴", "HbA1c 逾期 243 天，建議今日開立檢驗", ACCENT_RED),
    ("09:00", "王大明", "🟡", "目前用藥 12 種，注意多重用藥", ACCENT_ORANGE),
    ("09:15", "李小華", "🟢", "所有指標正常", ACCENT_GREEN),
    ("09:30", "張美玲", "🔴", "流感疫苗未接種（糖尿病高風險族群）", ACCENT_RED),
    ("09:30", "張美玲", "🔴", "血壓未追蹤超過 3 個月", ACCENT_RED),
    ("09:45", "陳志明", "🟡", "預防性抗生素已開立 2 天，注意第 3 天評估", ACCENT_ORANGE),
    ("10:00", "林美惠", "🟢", "所有指標正常", ACCENT_GREEN),
]

for i, (time, name, icon, msg, color) in enumerate(alerts):
    y = Inches(2.3 + i * 0.65)
    if i % 2 == 0:
        add_shape(slide, Inches(1.8), y, Inches(9.7), Inches(0.6), LIGHT_GRAY)
    add_text(slide, Inches(2.0), y + Inches(0.05), Inches(0.8), Inches(0.5),
             time, 14, GRAY, False)
    add_text(slide, Inches(2.8), y + Inches(0.05), Inches(1.2), Inches(0.5),
             name, 14, DARK_TEXT, True)
    add_text(slide, Inches(4.0), y + Inches(0.05), Inches(0.4), Inches(0.5),
             icon, 14, color, False)
    add_text(slide, Inches(4.5), y + Inches(0.05), Inches(7.0), Inches(0.5),
             msg, 14, DARK_TEXT, False)

add_text(slide, Inches(1.8), Inches(6.6), Inches(9.5), Inches(0.4),
         "💡 不需要連接 HIS，只需要 FHIR Server 中的掛號清單與病歷資料", 13, GRAY, False, PP_ALIGN.CENTER)


# ========================================================
# Slide 7: 個人品質指標 - 分子分母狀態解釋
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "個人品質指標：分母、分子、狀態怎麼看？", 30, DARK_BLUE, True)

# 三個概念卡片
# 分母
add_rounded_shape(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(2.2), LIGHT_BLUE, BRAND_BLUE)
add_text(slide, Inches(0.7), Inches(1.4), Inches(3.4), Inches(0.5),
         "分母 = 該不該管他？", 20, BRAND_BLUE, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(0.9), Inches(2.0), Inches(3.0), Inches(0.03), BRAND_BLUE)
add_multi_text(slide, Inches(0.9), Inches(2.2), Inches(3.2), Inches(1.2), [
    ("✅ 有糖尿病", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("→ HbA1c 指標管得到他", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("— 是男性", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("→ 剖腹產指標管不到他", 14, GRAY, False, PP_ALIGN.LEFT),
])

# 分子
add_rounded_shape(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(2.2), RGBColor(0xF0, 0xFF, 0xF0), ACCENT_GREEN)
add_text(slide, Inches(4.9), Inches(1.4), Inches(3.4), Inches(0.5),
         "分子 = 該做的做了沒？", 20, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(5.1), Inches(2.0), Inches(3.0), Inches(0.03), ACCENT_GREEN)
add_multi_text(slide, Inches(5.1), Inches(2.2), Inches(3.2), Inches(1.2), [
    ("✅ 有開連續處方", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("→ 做了！達標", 14, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("❌ 254 天沒做 HbA1c", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("→ 沒做！未達標", 14, ACCENT_RED, False, PP_ALIGN.LEFT),
])

# 狀態
add_rounded_shape(slide, Inches(8.9), Inches(1.3), Inches(3.8), Inches(2.2), RGBColor(0xFF, 0xF3, 0xE0), ACCENT_ORANGE)
add_text(slide, Inches(9.1), Inches(1.4), Inches(3.4), Inches(0.5),
         "狀態 = 要不要提醒？", 20, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(9.3), Inches(2.0), Inches(3.0), Inches(0.03), ACCENT_ORANGE)
add_multi_text(slide, Inches(9.3), Inches(2.2), Inches(3.2), Inches(1.2), [
    ("🟢 該管 + 做了 = 沒事", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("", 4, GRAY, False, PP_ALIGN.LEFT),
    ("🔴 該管 + 沒做 = 警示！", 15, ACCENT_RED, True, PP_ALIGN.LEFT),
    ("", 4, GRAY, False, PP_ALIGN.LEFT),
    ("⚪ 不適用 = 跳過", 15, GRAY, False, PP_ALIGN.LEFT),
])

# 實際範例表格
add_text(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(0.5),
         "範例：王大明　M/65　糖尿病+高血壓", 20, DARK_BLUE, True)

example_data = [
    ["指標", "分母\n（該不該管？）", "分子\n（做了沒？）", "狀態", "白話說明"],
    ["07 HbA1c", "✅ 有糖尿病\n所以要管", "❌ 254 天沒做\n沒做！", "🔴", "提醒醫師：\n該做 HbA1c 了！"],
    ["04 慢性處方", "✅ 有慢性病\n所以要管", "✅ 已開連續處方\n做了！", "🟢", "很好，\n不用提醒"],
    ["05 多重用藥", "✅ 有在吃藥\n所以要管", "❌ 吃了 12 種\n太多了！", "🔴", "提醒醫師：\n藥太多，該檢視"],
    ["08 血壓追蹤", "✅ 有高血壓\n所以要管", "✅ 上月已量\n做了！", "🟢", "很好，\n不用提醒"],
    ["11 剖腹產", "— 男性\n管不到", "—", "⚪", "跟他無關\n跳過"],
    ["12 抗生素", "— 沒做手術\n管不到", "—", "⚪", "跟他無關\n跳過"],
]

add_table(slide, Inches(0.3), Inches(4.3), Inches(12.7), Inches(3.0), example_data,
          col_widths=[Inches(1.8), Inches(2.8), Inches(2.8), Inches(1.0), Inches(4.3)])

# 底部總結
add_rounded_shape(slide, Inches(2.5), Inches(7.05), Inches(8.3), Inches(0.35), BRAND_BLUE)
add_text(slide, Inches(2.7), Inches(7.07), Inches(7.9), Inches(0.3),
         "適用指標：4/39　　達標：2/4　　待改善：2/4　→　醫師只需關注 🔴 的項目", 14, WHITE, True, PP_ALIGN.CENTER)


# ========================================================
# Slide 8: 資料轉換流程
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "資料轉換：SAS EG 表格 → FHIR 格式", 30, DARK_BLUE, True)

# 步驟流程
steps = [
    ("Step 1", "貴團隊提供\nSAS EG 欄位清單", "告訴我們有哪些表、\n哪些欄位", BRAND_BLUE),
    ("Step 2", "共同填寫\n欄位對照表", "SAS 欄位 → FHIR 欄位\n一起定義 mapping", ACCENT_GREEN),
    ("Step 3", "我們撰寫\n轉換腳本", "Python 自動讀 CSV\n寫入 FHIR Server", ACCENT_ORANGE),
    ("Step 4", "每日批次\n自動轉換", "SAS EG 匯出 CSV\n腳本自動匯入 FHIR", ACCENT_PURPLE),
]

for i, (step, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 3.2)
    # 步驟圓
    add_rounded_shape(slide, x, Inches(1.5), Inches(2.8), Inches(1.2), color)
    add_text(slide, x + Inches(0.1), Inches(1.55), Inches(2.6), Inches(0.4),
             step, 14, WHITE, False, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(1.85), Inches(2.6), Inches(0.8),
             title, 16, WHITE, True, PP_ALIGN.CENTER)

    # 箭頭
    if i < 3:
        add_text(slide, x + Inches(2.8), Inches(1.8), Inches(0.4), Inches(0.5),
                 "→", 28, GRAY, True, PP_ALIGN.CENTER)

    # 說明
    add_multi_text(slide, x + Inches(0.1), Inches(2.9), Inches(2.6), Inches(1.0), [
        (line, 13, DARK_TEXT, False, PP_ALIGN.CENTER) for line in desc.split('\n')
    ])

# 欄位對照表範例
add_text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5),
         "欄位對照範例", 20, DARK_BLUE, True)

mapping_data = [
    ["SAS EG 表", "SAS 欄位", "FHIR Resource", "FHIR 欄位"],
    ["病人基本", "病歷號", "Patient", "identifier.value"],
    ["門診紀錄", "就診日期", "Encounter", "period.start"],
    ["門診紀錄", "ICD 診斷碼", "Condition", "code.coding.code"],
    ["用藥紀錄", "健保碼", "MedicationRequest", "medication.coding.code"],
    ["檢驗報告", "HbA1c 數值", "Observation", "valueQuantity.value"],
]

add_table(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.3), mapping_data,
          col_widths=[Inches(2.5), Inches(2.5), Inches(3.5), Inches(3.8)])


# ========================================================
# Slide 8: 具體合作範例 - 指標交叉驗證
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "第一步：挑 3 個指標做交叉驗證", 30, DARK_BLUE, True)

verify_data = [
    ["品質指標", "SAS 計算方式", "CQL 計算方式", "預期驗證結果"],
    ["糖尿病 HbA1c\n檢測率", "SAS 撈 LAB_RESULTS\n計算有做 HbA1c 的比率", "CQL 查 Observation\ncode = LOINC 4548-4", "兩邊比率應一致\n差異 < 1%"],
    ["抗生素使用\n超過 3 天率", "SAS 撈 DRUG_ORDERS\n篩選抗生素 > 3 天", "CQL 查 MedicationRequest\n計算 duration > 3", "驗證收案條件\n是否完全一致"],
    ["剖腹產率", "SAS 撈手術紀錄\n篩選剖腹產", "CQL 查 Procedure\ncode = 健保處置碼", "確認處置碼\n對照正確性"],
]

add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.5), verify_data,
          col_widths=[Inches(2.3), Inches(3.3), Inches(3.4), Inches(3.3)])

add_rounded_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), RGBColor(0xFE, 0xF9, 0xE7), ACCENT_ORANGE)
add_multi_text(slide, Inches(1.2), Inches(5.4), Inches(11.0), Inches(1.5), [
    ("驗證的意義", 18, ACCENT_ORANGE, True, PP_ALIGN.LEFT),
    ("✅ 結果一致 → CQL 邏輯獲得實戰驗證，可信度大增", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("✅ 結果不一致 → 找出差異原因（收案條件？排除條件？時間？），雙方都受益", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("✅ 對醫院來說 → 同一指標有兩套獨立驗證，品質管理更可靠", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
])


# ========================================================
# Slide 9: 技術架構 - 完整串接圖
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "合作後的完整架構", 30, DARK_BLUE, True)

# 流程圖 - 用形狀模擬
flow_items = [
    (Inches(5.0), Inches(1.2), "HIS 資料", DARK_BLUE),
    (Inches(5.0), Inches(2.2), "數據中台 + SAS EG", BRAND_BLUE),
    (Inches(1.5), Inches(3.5), "SAS VA\n院內儀表板", BRAND_BLUE),
    (Inches(5.0), Inches(3.5), "FHIR Server\n（TW Core 格式）", ACCENT_GREEN),
    (Inches(8.5), Inches(3.5), "CQL 規則引擎\n（19+ 指標）", ACCENT_ORANGE),
    (Inches(1.5), Inches(5.2), "院長/品管師\n看報表", GRAY),
    (Inches(5.0), Inches(5.2), "品質儀表板\n（事後統計）", ACCENT_GREEN),
    (Inches(8.5), Inches(5.2), "即時警示 App\n（事前預警）", ACCENT_ORANGE),
    (Inches(5.0), Inches(6.3), "院際交換資料", ACCENT_PURPLE),
]

for (x, y, text, color) in flow_items:
    add_rounded_shape(slide, x, y, Inches(2.8), Inches(0.8), color)
    add_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.6), Inches(0.7),
             text, 13, WHITE, True, PP_ALIGN.CENTER)

# 箭頭文字
arrows = [
    (Inches(6.0), Inches(1.85), "↓"),
    (Inches(3.2), Inches(3.0), "↙"),
    (Inches(6.0), Inches(3.0), "↓"),
    (Inches(8.0), Inches(3.0), "↘"),
    (Inches(3.0), Inches(4.5), "↓"),
    (Inches(6.0), Inches(4.5), "↓"),
    (Inches(9.5), Inches(4.5), "↓"),
    (Inches(6.0), Inches(6.1), "↓"),
]

for (x, y, arrow) in arrows:
    add_text(slide, x, y, Inches(0.5), Inches(0.5), arrow, 22, GRAY, True, PP_ALIGN.CENTER)

# 標註
add_text(slide, Inches(3.5), Inches(3.1), Inches(2.0), Inches(0.3),
         "CSV 匯出", 11, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(7.0), Inches(3.1), Inches(2.0), Inches(0.3),
         "Python 轉換", 11, GRAY, False, PP_ALIGN.CENTER)


# ========================================================
# Slide 10: 對雙方的好處
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "合作對雙方的好處", 30, DARK_BLUE, True)

# 左：對 SAS 團隊
add_rounded_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), LIGHT_BLUE, BRAND_BLUE)
add_text(slide, Inches(1.0), Inches(1.4), Inches(5.1), Inches(0.5),
         "對貴團隊的好處", 22, BRAND_BLUE, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.03), BRAND_BLUE)

add_multi_text(slide, Inches(1.2), Inches(2.3), Inches(4.8), Inches(4.2), [
    ("1. 資料價值延伸", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   SAS EG 整理好的資料，不只用於 SAS VA", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   還能轉成國際標準，價值倍增", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("2. 即時功能加值", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   SAS VA 做不到的即時警示", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   透過合作可以做到", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("3. 跨院展示能力", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   院際交換是健保署未來方向", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   合作後直接具備此能力", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("4. 雙重驗證品質", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   兩套獨立系統交叉驗證", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   品管報告更有公信力", 14, GRAY, False, PP_ALIGN.LEFT),
])

# 右：對我們
add_rounded_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5), RGBColor(0xF0, 0xFF, 0xF0), ACCENT_GREEN)
add_text(slide, Inches(7.2), Inches(1.4), Inches(5.1), Inches(0.5),
         "對我們的好處", 22, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(7.4), Inches(2.0), Inches(4.8), Inches(0.03), ACCENT_GREEN)

add_multi_text(slide, Inches(7.4), Inches(2.3), Inches(4.8), Inches(4.2), [
    ("1. 真實資料來源", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   不用花時間處理原始 HIS 資料", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   貴團隊已經整理好了", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("2. 實戰驗證 CQL", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   CQL 規則對照 SAS 結果", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   獲得最有力的正確性證明", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("3. 落地應用場景", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   從「可以做」變成「正在用」", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   對投資人 / 評審有說服力", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("4. 合作publication", 17, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("   雙方共同發表", 14, GRAY, False, PP_ALIGN.LEFT),
    ("   SAS vs CQL 品質量測比較研究", 14, GRAY, False, PP_ALIGN.LEFT),
])


# ========================================================
# Slide 11: 資安與部署 - 單機 FHIR Server
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "資料保密：院內單機部署方案", 30, DARK_BLUE, True)

# 架構圖
add_rounded_shape(slide, Inches(1.5), Inches(1.3), Inches(10.3), Inches(3.2), LIGHT_BLUE, BRAND_BLUE)
add_text(slide, Inches(1.7), Inches(1.4), Inches(9.9), Inches(0.4),
         "醫院內部網路（完全不接外網）", 18, BRAND_BLUE, True, PP_ALIGN.CENTER)

# FHIR PC
add_rounded_shape(slide, Inches(2.5), Inches(2.0), Inches(3.0), Inches(2.0), BRAND_BLUE)
add_multi_text(slide, Inches(2.7), Inches(2.1), Inches(2.6), Inches(1.8), [
    ("FHIR Server PC", 16, WHITE, True, PP_ALIGN.CENTER),
    ("", 6, WHITE, False, PP_ALIGN.CENTER),
    ("HAPI FHIR（開源免費）", 13, WHITE, False, PP_ALIGN.CENTER),
    ("CQL 引擎", 13, WHITE, False, PP_ALIGN.CENTER),
    ("警示 App", 13, WHITE, False, PP_ALIGN.CENTER),
    ("品質儀表板", 13, WHITE, False, PP_ALIGN.CENTER),
])

# 連線箭頭
add_text(slide, Inches(5.7), Inches(2.7), Inches(1.5), Inches(0.5),
         "← 區網 →", 16, DARK_TEXT, True, PP_ALIGN.CENTER)

# 醫師電腦
add_rounded_shape(slide, Inches(7.5), Inches(2.0), Inches(3.0), Inches(2.0), ACCENT_GREEN)
add_multi_text(slide, Inches(7.7), Inches(2.1), Inches(2.6), Inches(1.8), [
    ("醫師 / 護理站電腦", 16, WHITE, True, PP_ALIGN.CENTER),
    ("", 6, WHITE, False, PP_ALIGN.CENTER),
    ("開瀏覽器即可使用", 13, WHITE, False, PP_ALIGN.CENTER),
    ("不需安裝任何軟體", 13, WHITE, False, PP_ALIGN.CENTER),
    ("Windows 防火牆", 13, WHITE, False, PP_ALIGN.CENTER),
    ("限制指定 IP 連入", 13, WHITE, False, PP_ALIGN.CENTER),
])

# 資安保證表
security_data = [
    ["資安疑慮", "解決方式"],
    ["資料外洩", "單機部署，完全不接外網"],
    ["未授權存取", "Windows 防火牆限制 IP，僅允許指定電腦"],
    ["傳輸安全", "FHIR Server 支援 HTTPS 加密"],
    ["個資保護", "轉換腳本自動去除姓名、身分證（去識別化）"],
    ["存取記錄", "HAPI FHIR 內建 Audit Log，可追溯所有操作"],
    ["IRB 審查", "品質改善通常免 IRB；若做研究再另行申請"],
]

add_table(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.7), security_data,
          col_widths=[Inches(3.5), Inches(8.8)])


# ========================================================
# Slide 12: 費用評估 - 零軟體費
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "費用評估：零軟體授權費", 30, DARK_BLUE, True)

# 軟體費用
add_text(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
         "軟體費用", 22, DARK_BLUE, True)

sw_data = [
    ["項目", "費用"],
    ["HAPI FHIR Server", "免費（開源）"],
    ["Java Runtime", "免費"],
    ["Docker Desktop", "免費（教育/個人）"],
    ["作業系統", "Windows / Linux 皆可"],
    ["資料庫", "H2 內建免費 / PostgreSQL 免費"],
    ["CQL 引擎", "免費（已開發完成）"],
    ["品質平台 + 警示 App", "免費（已開發完成）"],
    ["合計軟體費", "NT$ 0"],
]

add_table(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(3.5), sw_data,
          col_widths=[Inches(3.0), Inches(2.5)])

# 硬體需求
add_text(slide, Inches(7.0), Inches(1.2), Inches(5), Inches(0.5),
         "硬體需求（一台 PC 即可）", 22, DARK_BLUE, True)

hw_data = [
    ["規格", "最低需求", "建議配置"],
    ["CPU", "雙核", "四核 i5 以上"],
    ["記憶體", "4 GB", "8-16 GB"],
    ["硬碟", "20 GB", "100 GB SSD"],
    ["網路", "區域網路", "區域網路"],
    ["外網", "不需要", "不需要"],
]

add_table(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.5), hw_data,
          col_widths=[Inches(1.5), Inches(2.0), Inches(2.3)])

# 對比 SAS
add_rounded_shape(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.6), RGBColor(0xFE, 0xF9, 0xE7), ACCENT_ORANGE)
add_multi_text(slide, Inches(1.2), Inches(5.7), Inches(11.0), Inches(1.4), [
    ("對比：SAS 授權費 vs 我們的方案", 18, ACCENT_ORANGE, True, PP_ALIGN.LEFT),
    ("SAS VA 企業版授權：每年數十萬～數百萬元", 15, DARK_TEXT, False, PP_ALIGN.LEFT),
    ("FHIR-CQL 方案：軟體 NT$0 + 一台舊 PC = 立即上線", 15, ACCENT_GREEN, True, PP_ALIGN.LEFT),
    ("→ 省下的授權費可以用在更有價值的地方", 14, GRAY, False, PP_ALIGN.LEFT),
])


# ========================================================
# Slide 13: 可攜式部署 - USB 隨身碟帶著走
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), BRAND_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
         "可攜式部署：USB 帶著走，30 分鐘上線", 30, DARK_BLUE, True)

# USB 內容
add_rounded_shape(slide, Inches(0.8), Inches(1.3), Inches(5.0), Inches(5.3), LIGHT_BLUE, BRAND_BLUE)
add_text(slide, Inches(1.0), Inches(1.4), Inches(4.6), Inches(0.5),
         "USB 隨身碟內容", 20, BRAND_BLUE, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(1.2), Inches(2.0), Inches(4.2), Inches(0.03), BRAND_BLUE)

add_multi_text(slide, Inches(1.2), Inches(2.2), Inches(4.4), Inches(4.2), [
    ("📁 1_FHIR_Server/", 15, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("    docker-compose.yml", 13, GRAY, False, PP_ALIGN.LEFT),
    ("    hapi-fhir.jar（備用）", 13, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("📁 2_CQL_規則/", 15, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("    19+ 個品質指標 CQL 檔案", 13, GRAY, False, PP_ALIGN.LEFT),
    ("    編譯好的 ELM JSON", 13, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("📁 3_轉換腳本/", 15, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("    convert_sas_to_fhir.py", 13, GRAY, False, PP_ALIGN.LEFT),
    ("    mapping_config.json", 13, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("📁 4_警示App/", 15, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("    品質平台前端 + CQL 引擎", 13, GRAY, False, PP_ALIGN.LEFT),
    ("", 6, GRAY, False, PP_ALIGN.LEFT),
    ("📁 5_安裝指引/", 15, DARK_TEXT, True, PP_ALIGN.LEFT),
    ("    一頁安裝說明", 13, GRAY, False, PP_ALIGN.LEFT),
])

# 到院部署流程
add_rounded_shape(slide, Inches(6.5), Inches(1.3), Inches(6.0), Inches(5.3), RGBColor(0xF0, 0xFF, 0xF0), ACCENT_GREEN)
add_text(slide, Inches(6.7), Inches(1.4), Inches(5.6), Inches(0.5),
         "到院部署流程（30 分鐘）", 20, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(6.9), Inches(2.0), Inches(5.2), Inches(0.03), ACCENT_GREEN)

deploy_steps = [
    ("Step 1", "5 min", "安裝 Docker 或 Java"),
    ("Step 2", "2 min", "啟動 FHIR Server"),
    ("Step 3", "5 min", "拿到 SAS EG 的 CSV 匯出檔"),
    ("Step 4", "3 min", "跑轉換腳本，資料進 FHIR Server"),
    ("Step 5", "5 min", "開瀏覽器，打開警示 App"),
    ("Step 6", "10 min", "Demo 給醫師看"),
]

for i, (step, time_str, desc) in enumerate(deploy_steps):
    y = Inches(2.2 + i * 0.65)
    # 步驟
    add_rounded_shape(slide, Inches(6.9), y, Inches(1.0), Inches(0.5), ACCENT_GREEN)
    add_text(slide, Inches(6.9), y + Inches(0.05), Inches(1.0), Inches(0.4),
             step, 12, WHITE, True, PP_ALIGN.CENTER)
    # 時間
    add_text(slide, Inches(8.0), y + Inches(0.05), Inches(1.0), Inches(0.4),
             time_str, 13, ACCENT_GREEN, True, PP_ALIGN.LEFT)
    # 說明
    add_text(slide, Inches(9.1), y + Inches(0.05), Inches(3.2), Inches(0.4),
             desc, 14, DARK_TEXT, False, PP_ALIGN.LEFT)

# 底部重點
add_rounded_shape(slide, Inches(6.9), Inches(6.2), Inches(5.2), Inches(0.35), ACCENT_ORANGE)
add_text(slide, Inches(7.0), Inches(6.22), Inches(5.0), Inches(0.3),
         "程式碼完全不需修改，只有資料不同", 14, WHITE, True, PP_ALIGN.CENTER)


# ========================================================
# Slide 14: 下一步行動
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
         "下一步行動", 32, WHITE, True)

steps_next = [
    ("1", "本週", "貴團隊提供 SAS EG 的欄位清單\n（一份範例 CSV，10 筆資料即可）", ACCENT_BLUE),
    ("2", "下週", "共同填寫欄位對照表\n定義 SAS 欄位 → FHIR 欄位的 mapping", ACCENT_GREEN),
    ("3", "第 2-3 週", "我們撰寫 Python 轉換腳本\n並上傳測試資料到 FHIR Server", ACCENT_ORANGE),
    ("4", "第 3-4 週", "挑 3 個指標做交叉驗證\n（糖尿病 HbA1c / 抗生素 / 剖腹產率）", ACCENT_PURPLE),
    ("5", "第 5-6 週", "開發即時警示 App 原型\n護理站第二視窗展示", RGBColor(0xE8, 0x3E, 0x3E)),
]

for i, (num, timeline, desc, color) in enumerate(steps_next):
    y = Inches(1.5 + i * 1.1)
    # 數字圓圈
    add_rounded_shape(slide, Inches(1.0), y, Inches(0.7), Inches(0.7), color)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(0.7), Inches(0.5),
             num, 22, WHITE, True, PP_ALIGN.CENTER)
    # 時間
    add_text(slide, Inches(2.0), y + Inches(0.1), Inches(1.5), Inches(0.5),
             timeline, 16, color, True)
    # 說明
    add_multi_text(slide, Inches(3.8), y, Inches(8.5), Inches(0.9), [
        (line, 15, WHITE, False, PP_ALIGN.LEFT) for line in desc.split('\n')
    ])

add_rounded_shape(slide, Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.3), ACCENT_GREEN)
add_text(slide, Inches(1.2), Inches(6.95), Inches(11.0), Inches(0.3),
         "📌 第一步最簡單：只需要一份 10 筆資料的範例 CSV", 15, WHITE, True, PP_ALIGN.CENTER)


# ========================================================
# Slide 12: 結語
# ========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
         "你們把資料洗乾淨（最苦的活）", 30, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.0),
         "我們把資料標準化 + 智慧化（最難的活）", 30, ACCENT_BLUE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(1.0),
         "合在一起，就是從 HIS 到臨床決策的完整解決方案", 26, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(3), Inches(5.5), Inches(7.3), Inches(0.03), GRAY)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
         "報表讓院長看過去，警示讓醫師改現在", 22, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
         "謝謝　期待合作", 28, WHITE, True, PP_ALIGN.CENTER)


# ========================================================
# 儲存
# ========================================================
output_path = r"合作提案_FHIR-CQL_x_SAS團隊.pptx"
prs.save(output_path)
print(f"✅ 簡報已產生：{output_path}")
print(f"   共 {len(prs.slides)} 張投影片")
