"""
FHIR CQL 平台 — Pitch Deck V3（字少、重點、簡潔）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 色彩 ──
BG_DARK = RGBColor(0x0f, 0x17, 0x2a)
BG_LIGHT = RGBColor(0xf8, 0xfa, 0xfc)
BLUE = RGBColor(0x38, 0xbd, 0xf8)
GREEN = RGBColor(0x10, 0xb9, 0x81)
ORANGE = RGBColor(0xfb, 0x92, 0x3c)
RED = RGBColor(0xef, 0x44, 0x44)
PURPLE = RGBColor(0xa7, 0x8b, 0xfa)
TEAL = RGBColor(0x06, 0xb6, 0xd4)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x94, 0xa3, 0xb8)
DARK = RGBColor(0x1e, 0x29, 0x3b)
CARD = RGBColor(0x1e, 0x29, 0x3b)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, border=None, bw=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def rr(slide, l, t, w, h, fill, border=None, bw=0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = "Microsoft JhengHei"; p.alignment = align
    return tb


def multi(slide, l, t, w, h, lines, size=16, color=WHITE, spacing=1.5):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            tx, bl, cl = item, False, color
        else:
            tx = item[0]; bl = item[1] if len(item) > 1 else False; cl = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = tx; p.font.size = Pt(size); p.font.color.rgb = cl
        p.font.bold = bl; p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(size * (spacing - 1))
    return tb


# ═══════════════════════════════════════════════
# 1. 封面
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), BLUE)

txt(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
    "FHIR CQL Platform", 52, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.2), Inches(11), Inches(0.7),
    "軟硬整合 ╳ 伺服器廠商合作提案", 26, BLUE, False, PP_ALIGN.CENTER)

rect(s, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.03), BLUE)

txt(s, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
    "台灣唯一 CQL 執行平台  ·  得獎產品  ·  50+ 品質指標模組", 16, GRAY, False, PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
    "2026.04", 14, GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 2. 問題
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), RED)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "問題", 40, RED, True)

# 三個大痛點
pains = [
    ("CQL 寫好了\n沒地方跑", "衛福部推 FHIR+CQL\n但全台零執行平台"),
    ("醫院不開放\n資料撈取", "不可能讓每個 CQL 創作者\n連進 FHIR Server"),
    ("品質申報\n仍靠人工", "Excel + 人工抽審\n475 家醫院都這樣做"),
]
for i, (big, small) in enumerate(pains):
    x = Inches(1 + i * 4)
    rr(s, x, Inches(2.0), Inches(3.5), Inches(3.5), RGBColor(0x1a, 0x24, 0x3b), RED, 1.5)
    multi(s, x + Inches(0.4), Inches(2.4), Inches(2.7), Inches(1.5),
           big.split('\n'), 26, WHITE, 1.3)
    multi(s, x + Inches(0.4), Inches(4.0), Inches(2.7), Inches(1.2),
           small.split('\n'), 14, GRAY, 1.5)

txt(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
    "每年健保支出 8,700 億+，品質量測卻停留在紙本時代", 18, RED, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 3. 解法
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), GREEN)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "解法", 40, GREEN, True)

txt(s, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
    "CQL Marketplace ＋ 院內執行引擎", 28, WHITE, True, PP_ALIGN.CENTER)

# 左：Marketplace
rr(s, Inches(1), Inches(2.8), Inches(5), Inches(3.5), RGBColor(0x1a, 0x24, 0x3b), GREEN, 1.5)
txt(s, Inches(1.5), Inches(3.0), Inches(4), Inches(0.5), "☁️  CQL Marketplace", 20, GREEN, True)
mp_items = ["CQL 上架 / 審核 / 版本管理", "創作者生態系", "授權計費"]
for j, item in enumerate(mp_items):
    txt(s, Inches(1.5), Inches(3.7 + j * 0.6), Inches(4), Inches(0.4), f"›  {item}", 16, GRAY)

# 箭頭
txt(s, Inches(6.2), Inches(4.0), Inches(1), Inches(0.8), "⟹", 40, BLUE, True, PP_ALIGN.CENTER)

# 右：院內引擎
rr(s, Inches(7.3), Inches(2.8), Inches(5), Inches(3.5), RGBColor(0x1a, 0x24, 0x3b), BLUE, 1.5)
txt(s, Inches(7.8), Inches(3.0), Inches(4), Inches(0.5), "🏥  院內 Appliance", 20, BLUE, True)
eng_items = ["CQL 執行引擎 + FHIR 適配", "資料永遠不出院", "統計結果回傳"]
for j, item in enumerate(eng_items):
    txt(s, Inches(7.8), Inches(3.7 + j * 0.6), Inches(4), Inches(0.4), f"›  {item}", 16, GRAY)

txt(s, Inches(1), Inches(6.6), Inches(11), Inches(0.4),
    "資料不出院  ·  CQL 進去跑  ·  結果出來", 16, GREEN, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 4. 合作模式（一頁看完）
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), PURPLE)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "合作", 40, PURPLE, True)

# 我方
rr(s, Inches(1), Inches(2.0), Inches(5), Inches(3.0), RGBColor(0x1a, 0x24, 0x3b), PURPLE, 1.5)
txt(s, Inches(1.5), Inches(2.2), Inches(4), Inches(0.5), "🧠  我方", 22, PURPLE, True)
my = ["CQL 執行引擎 & 50+ 指標", "Marketplace 平台", "品質報表 & 儀表板", "技術支援 & 更新"]
for j, item in enumerate(my):
    txt(s, Inches(1.5), Inches(3.0 + j * 0.5), Inches(4), Inches(0.4), item, 16, WHITE)

# 貴方
rr(s, Inches(7.3), Inches(2.0), Inches(5), Inches(3.0), RGBColor(0x1a, 0x24, 0x3b), BLUE, 1.5)
txt(s, Inches(7.8), Inches(2.2), Inches(4), Inches(0.5), "🖥️  貴方", 22, BLUE, True)
their = ["伺服器硬體", "醫院通路 & 業務團隊", "機房維運 & SLA", "資安合規"]
for j, item in enumerate(their):
    txt(s, Inches(7.8), Inches(3.0 + j * 0.5), Inches(4), Inches(0.4), item, 16, WHITE)

# 底部整合
rect(s, Inches(1), Inches(5.5), Inches(11.3), Inches(0.8), PURPLE)
txt(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.6),
    "醫院買一台 → 開箱即用 → 連上 Marketplace → 即刻運行", 20, WHITE, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 5. 市場
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), ORANGE)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "市場", 40, ORANGE, True)

# 大數字
nums = [
    ("8,700 億+", "健保年支出", ORANGE),
    ("475 家", "全台醫院", BLUE),
    ("19 家", "醫學中心\n（優先目標）", GREEN),
    ("99.9%", "健保覆蓋率", PURPLE),
]
for i, (n, label, clr) in enumerate(nums):
    x = Inches(1 + i * 3.1)
    rr(s, x, Inches(2.0), Inches(2.7), Inches(1.8), RGBColor(0x1a, 0x24, 0x3b), clr, 1.5)
    txt(s, x, Inches(2.2), Inches(2.7), Inches(0.8), n, 30, clr, True, PP_ALIGN.CENTER)
    multi(s, x + Inches(0.2), Inches(3.0), Inches(2.3), Inches(0.7),
           label.split('\n'), 13, GRAY, 1.3)

# 營收模型（極簡版）
txt(s, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
    "3 年保守目標", 20, WHITE, True)

targets = [
    ("醫學中心  6 家 × 60 萬", "360 萬"),
    ("區域醫院  13 家 × 40 萬", "520 萬"),
    ("地區醫院  37 家 × 20 萬", "740 萬"),
]
for j, (desc, rev) in enumerate(targets):
    y = Inches(5.0 + j * 0.55)
    txt(s, Inches(1.5), y, Inches(6), Inches(0.4), desc, 16, GRAY)
    txt(s, Inches(8), y, Inches(2), Inches(0.4), rev, 16, WHITE, True, PP_ALIGN.RIGHT)

rect(s, Inches(1.5), Inches(6.7), Inches(8.5), Inches(0.03), ORANGE)
txt(s, Inches(1.5), Inches(6.8), Inches(6), Inches(0.4), "合計年營收", 16, GRAY)
txt(s, Inches(8), Inches(6.8), Inches(2), Inches(0.4), "1,620 萬+", 20, ORANGE, True, PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════
# 6. 競爭
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), TEAL)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "競爭", 40, TEAL, True)

txt(s, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
    "台灣 CQL 執行平台數量：0", 28, RED, True, PP_ALIGN.CENTER)

# 簡潔比較
comps = [
    ("台灣 HIS 廠商", "有通路\n無 CQL、無 FHIR 品質量測", GRAY),
    ("國際 EHR\n(Epic, Oracle)", "有 CQL\n不懂台灣法規、編碼、指標", GRAY),
    ("Excel + 人工", "現行做法\n耗時、易錯、不可追溯", GRAY),
    ("我們", "CQL ✅  FHIR ✅  台灣 ✅\n50+ 指標 ✅  得獎 ✅", GREEN),
]
for i, (name, desc, clr) in enumerate(comps):
    x = Inches(1 + i * 3.1)
    is_us = (i == 3)
    border = GREEN if is_us else RGBColor(0x2a, 0x34, 0x4b)
    bw = 2.5 if is_us else 1
    rr(s, x, Inches(2.8), Inches(2.7), Inches(2.5), RGBColor(0x1a, 0x24, 0x3b), border, bw)
    multi(s, x + Inches(0.3), Inches(3.0), Inches(2.1), Inches(0.8),
           name.split('\n'), 16, WHITE if is_us else GRAY, 1.2)
    multi(s, x + Inches(0.3), Inches(3.9), Inches(2.1), Inches(1.2),
           desc.split('\n'), 13, GREEN if is_us else GRAY, 1.5)

# 壁壘
txt(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
    "護城河", 20, WHITE, True)
moats = ["先行者優勢（台灣零競爭）", "在地化壁壘（法規 + 編碼 + 指標定義）", "網路效應（CQL 越多 → 平台越有價值）"]
for j, m in enumerate(moats):
    txt(s, Inches(1 + j * 4), Inches(6.4), Inches(3.5), Inches(0.5), f"▸ {m}", 14, TEAL)


# ═══════════════════════════════════════════════
# 7. 商業模式
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), ORANGE)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "獲利", 40, ORANGE, True)

streams = [
    ("💰", "平台授權費", "醫院年費制", ORANGE),
    ("📦", "CQL 模組費", "按模組訂閱\n平台抽成", BLUE),
    ("🖥️", "硬體銷售分潤", "伺服器預裝\n合作分潤", GREEN),
    ("🔧", "客製化 & 維護", "專案計費\n年度合約", PURPLE),
]
for i, (icon, title, desc, clr) in enumerate(streams):
    x = Inches(1 + i * 3.1)
    rr(s, x, Inches(2.0), Inches(2.7), Inches(2.8), RGBColor(0x1a, 0x24, 0x3b), clr, 1.5)
    txt(s, x, Inches(2.2), Inches(2.7), Inches(0.7), icon, 36, WHITE, False, PP_ALIGN.CENTER)
    txt(s, x, Inches(2.9), Inches(2.7), Inches(0.5), title, 18, clr, True, PP_ALIGN.CENTER)
    multi(s, x + Inches(0.3), Inches(3.5), Inches(2.1), Inches(1.0),
           desc.split('\n'), 14, GRAY, 1.5)

# 飛輪
txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    "飛輪效應", 20, WHITE, True)
rect(s, Inches(1), Inches(6.1), Inches(11.3), Inches(0.8), RGBColor(0x1a, 0x24, 0x3b), ORANGE, 1.5)
txt(s, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.6),
    "更多 CQL 上架 → 更多醫院想用 → 更多伺服器賣出 → 更多創作者加入 → ♻️", 16, ORANGE, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 8. 時程
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), BLUE)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "時程", 40, BLUE, True)

phases = [
    ("P1", "POC", "1-2 月", ["Demo 機整合", "1-2 家醫院驗證"], BLUE),
    ("P2", "產品化", "3-4 月", ["Marketplace MVP", "定價 & MOU"], GREEN),
    ("P3", "擴展", "5-12 月", ["醫學中心導入", "區域醫院推廣"], ORANGE),
]
for i, (ph, title, dur, items, clr) in enumerate(phases):
    x = Inches(1 + i * 4.1)
    rr(s, x, Inches(2.0), Inches(3.6), Inches(4.2), RGBColor(0x1a, 0x24, 0x3b), clr, 1.5)
    # Phase badge
    rect(s, x + Inches(0.3), Inches(2.3), Inches(1), Inches(0.45), clr)
    txt(s, x + Inches(0.3), Inches(2.32), Inches(1), Inches(0.4), ph, 16, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(1.5), Inches(2.32), Inches(1.8), Inches(0.4), dur, 14, GRAY, False, PP_ALIGN.RIGHT)
    # Title
    txt(s, x + Inches(0.3), Inches(3.0), Inches(3), Inches(0.5), title, 24, WHITE, True)
    # Items
    for j, item in enumerate(items):
        txt(s, x + Inches(0.3), Inches(3.8 + j * 0.6), Inches(3), Inches(0.4), f"▸ {item}", 16, GRAY)


# ═══════════════════════════════════════════════
# 9. 為什麼是我們
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), GREEN)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "為什麼是我們", 40, GREEN, True)

reasons = [
    ("🏆", "得獎產品", "不是 PPT，是可運行的平台"),
    ("📊", "50+ CQL", "疾管 / 國健 / ESG / 品質四大類"),
    ("🥇", "唯一玩家", "台灣零 CQL 執行平台競爭者"),
    ("🔌", "多源接入", "FHIR / SQL / CSV 都能接"),
    ("📜", "專利申請中", "長庚共同發明，技術有背書"),
    ("🚀", "政策順風", "衛福部力推 FHIR，時機剛好"),
]
for i, (icon, title, desc) in enumerate(reasons):
    row = i // 3; col = i % 3
    x = Inches(1 + col * 4.1); y = Inches(2.0 + row * 2.5)
    rr(s, x, y, Inches(3.6), Inches(2.0), RGBColor(0x1a, 0x24, 0x3b), GREEN, 1)
    txt(s, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.6), icon, 28, WHITE)
    txt(s, x + Inches(1.1), y + Inches(0.25), Inches(2.2), Inches(0.4), title, 18, WHITE, True)
    txt(s, x + Inches(1.1), y + Inches(0.75), Inches(2.2), Inches(0.8), desc, 13, GRAY)


# ═══════════════════════════════════════════════
# 10. CDS Hooks 擴展
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), TEAL)

txt(s, Inches(1), Inches(0.8), Inches(4), Inches(0.8),
    "未來", 40, TEAL, True)

txt(s, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
    "同一平台，兩條產品線", 26, WHITE, True, PP_ALIGN.CENTER)

# 左
rr(s, Inches(1), Inches(2.8), Inches(5), Inches(2.5), RGBColor(0x1a, 0x24, 0x3b), GREEN, 1.5)
txt(s, Inches(1.5), Inches(3.0), Inches(4), Inches(0.5), "現在：品質量測", 20, GREEN, True)
txt(s, Inches(1.5), Inches(3.6), Inches(4), Inches(0.4), "批次 · 群體 · 報表", 16, GRAY)
txt(s, Inches(1.5), Inches(4.2), Inches(4), Inches(0.4), "「上一季抗生素使用率？」", 14, WHITE)

# 右
rr(s, Inches(7.3), Inches(2.8), Inches(5), Inches(2.5), RGBColor(0x1a, 0x24, 0x3b), TEAL, 1.5)
txt(s, Inches(7.8), Inches(3.0), Inches(4), Inches(0.5), "未來：CDS Hooks 即時決策", 20, TEAL, True)
txt(s, Inches(7.8), Inches(3.6), Inches(4), Inches(0.4), "即時 · 個人 · 推送", 16, GRAY)
txt(s, Inches(7.8), Inches(4.2), Inches(4), Inches(0.4), "「這個病人現在該不該開這藥？」", 14, WHITE)

# 加號
txt(s, Inches(6.2), Inches(3.5), Inches(0.8), Inches(0.8), "＋", 40, BLUE, True, PP_ALIGN.CENTER)

# 流程
txt(s, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
    "醫生開藥 → EHR 觸發 Hook → CQL 引擎判斷 → 回傳建議卡片", 16, TEAL, False, PP_ALIGN.CENTER)

rect(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.6), TEAL)
txt(s, Inches(1.5), Inches(6.55), Inches(10.3), Inches(0.5),
    "Marketplace 賣品質 CQL（現在）＋ 賣決策 CQL（未來）＝ 雙倍市場", 16, WHITE, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 11. 下一步
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, BG_DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), BLUE)

txt(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
    "下一步", 48, WHITE, True, PP_ALIGN.CENTER)

rect(s, Inches(5.5), Inches(2.6), Inches(2.3), Inches(0.03), BLUE)

steps = [
    ("1", "技術對接會議"),
    ("2", "POC Demo 機"),
    ("3", "簽 MOU"),
    ("4", "用成果說話"),
]
for i, (num, label) in enumerate(steps):
    x = Inches(1.5 + i * 2.8)
    rr(s, x, Inches(3.3), Inches(2.3), Inches(1.5), RGBColor(0x1a, 0x24, 0x3b), BLUE, 1.5)
    txt(s, x, Inches(3.4), Inches(2.3), Inches(0.6), num, 30, BLUE, True, PP_ALIGN.CENTER)
    txt(s, x, Inches(4.0), Inches(2.3), Inches(0.5), label, 16, WHITE, False, PP_ALIGN.CENTER)
    if i < 3:
        txt(s, x + Inches(2.3), Inches(3.8), Inches(0.5), Inches(0.5), "→", 20, BLUE, True, PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
    "「每家醫院一台，裡面跑我們的平台」", 26, BLUE, True, PP_ALIGN.CENTER)


# ── 儲存 ──
out = os.path.join(os.path.dirname(__file__), "FHIR_CQL_合作提案_V3.pptx")
prs.save(out)
print(f"✅ PPT 已生成：{out}")
